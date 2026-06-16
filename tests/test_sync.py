"""
test_sync.py — 手編集同期（sync.py）のテスト。

生成→編集→検出→反映 のループが壊れないことを保証する。
"""
import io
import pytest
from pptx import Presentation

import slidegen  # 型登録
from slidegen.parser import parse
from slidegen.render import build
from slidegen import sync


SRC = """\
slide prep
  headline "テスト見出し"
  col
    "最初の本文アルファ"
  col
    "二番目の本文ベータ"

---

slide kpt
  headline "ふりかえり見出し"
  col
    "続けることX"
  col
    "課題Y"
  col
    "試すことZ"
"""


def _render_to_pptx(src):
    prs = build(parse(src))
    return prs


def _edit_text(prs, find, replace):
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and find in sh.text_frame.text:
                sh.text_frame.paragraphs[0].runs[0].text = replace
                return True
    return False


def test_no_change_detected():
    """編集していなければ差分は出ない。"""
    prs = _render_to_pptx(SRC)
    prs.save("/tmp/_t_nochange.pptx")
    diff = sync.compute_diff(SRC, "/tmp/_t_nochange.pptx")
    total = sum(len(e["changes"]) for e in diff)
    assert total == 0


def test_single_change_detected_and_applied():
    """1箇所の編集を検出し、.slideに反映できる。"""
    prs = _render_to_pptx(SRC)
    assert _edit_text(prs, "最初の本文アルファ", "編集後アルファ")
    prs.save("/tmp/_t_single.pptx")

    diff = sync.compute_diff(SRC, "/tmp/_t_single.pptx")
    changes = [c for e in diff for c in e["changes"]]
    assert ("最初の本文アルファ", "編集後アルファ") in changes

    new_src, applied = sync.apply_changes(SRC, diff)
    assert applied == 1
    assert "編集後アルファ" in new_src
    assert "最初の本文アルファ" not in new_src

    # 反映後を再生成すると編集が活きている
    prs2 = _render_to_pptx(new_src)
    found = any("編集後アルファ" in sh.text_frame.text
                for s in prs2.slides for sh in s.shapes if sh.has_text_frame)
    assert found


def test_multi_slide_changes():
    """複数スライドの編集を検出できる。"""
    prs = _render_to_pptx(SRC)
    _edit_text(prs, "二番目の本文ベータ", "編集ベータ")
    _edit_text(prs, "ふりかえり見出し", "新ふりかえり見出し")
    prs.save("/tmp/_t_multi.pptx")

    diff = sync.compute_diff(SRC, "/tmp/_t_multi.pptx")
    changes = [c for e in diff for c in e["changes"]]
    assert ("二番目の本文ベータ", "編集ベータ") in changes
    assert ("ふりかえり見出し", "新ふりかえり見出し") in changes

    new_src, applied = sync.apply_changes(SRC, diff)
    assert applied == 2
    assert "編集ベータ" in new_src and "新ふりかえり見出し" in new_src


def test_apply_only_touches_changed_text():
    """変更のない文言は書き換えない。"""
    prs = _render_to_pptx(SRC)
    _edit_text(prs, "課題Y", "新しい課題")
    prs.save("/tmp/_t_safe.pptx")

    diff = sync.compute_diff(SRC, "/tmp/_t_safe.pptx")
    new_src, applied = sync.apply_changes(SRC, diff)
    # 他の文言は保持
    assert "続けることX" in new_src
    assert "試すことZ" in new_src
    assert "新しい課題" in new_src
