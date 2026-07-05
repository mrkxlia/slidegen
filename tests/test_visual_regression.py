"""
test_visual_regression.py — 第2層の自動化：全型の図形ツリー・スナップショット回帰テスト（案B）。

このテストは各型を最小入力でレンダし、図形ツリーの「決定的な部分集合」（種別・座標・塗り/線の色・
テキスト・フォント、表と本物チャートの構造）を正規化した JSON スナップショットと比較する。
pptx→画像のラスタライズは行わない（純Python・LibreOffice 不要・環境非依存）。
golden（tests/__snapshots__/visual_regression.json）はリポジトリにコミットされ、Git diff で
「どの型のどの属性が変わったか」が読める。意図した変更のときは下記コマンドで golden を更新する:

    SLIDEGEN_UPDATE_SNAPSHOTS=1 uv run --extra dev python -m pytest tests/test_visual_regression.py -q

【案A（画像スナップショット）へのフォールバック方針】
本テスト（案B）は図形ツリーの比較なので、図形ツリーは不変でも崩れる回帰
（描画順による重なり・微細なレイアウト崩れ 等）を取りこぼす可能性がある。また稀に脆すぎて
運用が回らなくなることもある。その場合は案A＝画像スナップショットへ移行する:
`tools/visual.py` の `slide_to_images()`（LibreOffice + pdftoppm）で代表型を画像化し、golden 画像と
ピクセル/知覚ハッシュで比較する。LibreOffice 依存で CI が重くなるため、案A採用時は対象を重要型に絞ること。
"""
from __future__ import annotations
import json
import os
import pathlib

import pytest
from pptx.util import Length

import slidegen  # 全 render_* を登録
from slidegen.render import RENDERERS, build
from slidegen.parser import Slide, Block

SNAP_DIR = pathlib.Path(__file__).resolve().parent / "__snapshots__"
SNAP_FILE = SNAP_DIR / "visual_regression.json"
_UPDATE = os.environ.get("SLIDEGEN_UPDATE_SNAPSHOTS") == "1"

# チャート系はカテゴリ＋数値が要るので、意味のあるスナップショットになるよう専用入力を与える。
_CHART_TYPES = ("bar_chart", "clustered_bar", "bar_horizontal",
                "line_chart", "stacked_bar", "stacked_100_bar")


def _minimal_slide(t: str) -> Slide:
    """全型に通用する最小入力（headline/kicker/foot + 2 ブロック）。"""
    return Slide(
        type=t,
        props={"headline": "見出しテスト", "kicker": "カテゴリ", "foot": "脚注テスト"},
        blocks=[
            Block(title="項目A", lines=["本文の一行目", "本文の二行目"]),
            Block(title="項目B", highlight=True, lines=["強調された項目"]),
        ],
    )


def _chart_slide(t: str) -> Slide:
    return Slide(
        type=t,
        props={"headline": "四半期売上", "categories_list": ["Q1", "Q2", "Q3", "Q4"],
               "unit": "百万円"},
        blocks=[
            Block(title="売上", lines=["120", "150", "135", "180"]),
            Block(title="目標", lines=["100", "140", "160", "170"]),
        ],
    )


def _input_for(t: str) -> Slide:
    return _chart_slide(t) if t in _CHART_TYPES else _minimal_slide(t)


# --- 図形ツリーの正規化シリアライズ（決定的な部分集合のみ） --------------------

def _hex(rgb) -> str:
    return str(rgb).upper()


def _solid_color(color_format) -> str | None:
    """SOLID かつ RGB のときだけ hex を返す（theme 色/scheme 色は None）。"""
    try:
        if color_format.type is None:
            return None
        return _hex(color_format.rgb)
    except Exception:
        return None


def _fill_of(shape) -> str | None:
    try:
        if shape.fill.type is None:
            return None
        return _solid_color(shape.fill.fore_color)
    except Exception:
        return None


def _line_of(shape) -> str | None:
    try:
        return _solid_color(shape.line.color)
    except Exception:
        return None


def _text_of(shape) -> list | None:
    if not getattr(shape, "has_text_frame", False):
        return None
    paras = []
    for p in shape.text_frame.paragraphs:
        runs = []
        for r in p.runs:
            size = r.font.size
            runs.append({
                "t": r.text,
                "font": r.font.name,
                "pt": round(Length(size).pt, 2) if size is not None else None,
                "bold": r.font.bold,
                "color": _solid_color(r.font.color),
            })
        if runs:
            paras.append(runs)
    return paras or None


def _table_of(shape) -> dict | None:
    if not getattr(shape, "has_table", False):
        return None
    tbl = shape.table
    rows = [[cell.text for cell in row.cells] for row in tbl.rows]
    return {"rows": len(rows), "cols": len(rows[0]) if rows else 0, "cells": rows}


def _chart_of(shape) -> dict | None:
    if not getattr(shape, "has_chart", False):
        return None
    ch = shape.chart
    try:
        plot0 = ch.plots[0]
        cats = list(plot0.categories)
    except Exception:
        cats = []
    series = []
    for s in ch.series:
        try:
            series.append({"name": s.name, "values": list(s.values)})
        except Exception:
            series.append({"name": None, "values": []})
    return {"chart_type": str(ch.chart_type), "categories": cats, "series": series}


def _shape_dict(shape) -> dict:
    st = shape.shape_type
    d = {
        "type": st.name if st is not None else None,
        "l": int(shape.left) if shape.left is not None else None,
        "t": int(shape.top) if shape.top is not None else None,
        "w": int(shape.width) if shape.width is not None else None,
        "h": int(shape.height) if shape.height is not None else None,
    }
    for key, val in (("fill", _fill_of(shape)), ("line", _line_of(shape)),
                     ("text", _text_of(shape)), ("table", _table_of(shape)),
                     ("chart", _chart_of(shape))):
        if val is not None:
            d[key] = val
    return d


def _serialize_shapes(slide) -> list:
    return [_shape_dict(sh) for sh in slide.shapes]


def _build_current() -> dict:
    raw = {}
    for t in sorted(RENDERERS):
        prs = build([_input_for(t)])
        raw[t] = _serialize_shapes(prs.slides[0])
    # JSON round-trip で GOLDEN と同じ表現（tuple→list 等）に正規化して比較を安定させる。
    return json.loads(json.dumps(raw, ensure_ascii=False))


CURRENT = _build_current()

if _UPDATE:
    SNAP_DIR.mkdir(parents=True, exist_ok=True)
    SNAP_FILE.write_text(
        json.dumps(CURRENT, ensure_ascii=False, sort_keys=True, indent=2) + "\n",
        encoding="utf-8",
    )
    GOLDEN = CURRENT
else:
    GOLDEN = json.loads(SNAP_FILE.read_text(encoding="utf-8")) if SNAP_FILE.exists() else {}


@pytest.mark.parametrize("t", sorted(RENDERERS))
def test_shape_tree_matches_snapshot(t):
    if not GOLDEN:
        pytest.fail(
            "golden スナップショットが未生成です。"
            "SLIDEGEN_UPDATE_SNAPSHOTS=1 で一度生成してコミットしてください。"
        )
    assert t in GOLDEN, f"{t} の golden が無い（SLIDEGEN_UPDATE_SNAPSHOTS=1 で更新）"
    assert CURRENT[t] == GOLDEN[t], (
        f"型 {t} の図形ツリーが golden と一致しません。意図した変更なら "
        f"SLIDEGEN_UPDATE_SNAPSHOTS=1 で golden を更新してください。"
    )


def test_snapshot_covers_all_renderers():
    """golden のキー集合が RENDERERS と一致（型の増減で golden 更新を強制＝ドリフト検知）。"""
    if not GOLDEN:
        pytest.skip("golden 未生成（初回は SLIDEGEN_UPDATE_SNAPSHOTS=1 で生成）")
    assert set(GOLDEN) == set(RENDERERS), (
        f"golden と RENDERERS の差分: "
        f"golden のみ={sorted(set(GOLDEN) - set(RENDERERS))} / "
        f"RENDERERS のみ={sorted(set(RENDERERS) - set(GOLDEN))}"
    )
