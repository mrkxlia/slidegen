"""
test_inspect_pptx.py — inspect_pptx（既存 pptx の構造抽出）のテスト。

デザイン取り込み（既存pptx→DSL再構成、LLM への入力）になる
inspect() / inspect_compact() を、slidegen 自身が生成した pptx で検証する。
純Python・LibreOffice 不要。
"""
import pathlib

import slidegen
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from slidegen.inspect_pptx import inspect, inspect_compact

REPO_ROOT = pathlib.Path(__file__).resolve().parent.parent
SAMPLE = (REPO_ROOT / "examples" / "sample.slide").read_text(encoding="utf-8")


def _sample_pptx(tmp_path) -> str:
    out = tmp_path / "deck.pptx"
    out.write_bytes(slidegen.render_to_bytes(SAMPLE))
    return str(out)


def test_inspect_returns_spec_fields(tmp_path):
    from slidegen.parser import parse

    data = inspect(_sample_pptx(tmp_path))
    assert len(data["slides"]) == len(parse(SAMPLE))
    assert len(data["slide_size_emu"]) == 2
    for spec in data["slides"]:
        assert spec["n_shapes"] == len(spec["shapes"])
        for s in spec["shapes"]:
            # 配置は%（0..100）で正規化されている。
            for k in ("x%", "y%", "w%", "h%"):
                if s[k] is not None:
                    assert -50 <= s[k] <= 150  # 意図的なはみ出し装飾を許容
    # 少なくとも1枚はフォント階層と面積パレットが取れる。
    assert any(spec["font_hierarchy"] for spec in data["slides"])
    assert any(spec["palette_by_area"] for spec in data["slides"])


def test_inspect_text_limit_expands_truncation(tmp_path):
    path = _sample_pptx(tmp_path)
    short = inspect(path, text_limit=5)
    long = inspect(path, text_limit=500)
    assert all(len(s["text"]) <= 5 for sp in short["slides"] for s in sp["shapes"])
    joined = lambda d: sum(len(s["text"]) for sp in d["slides"] for s in sp["shapes"])
    assert joined(long) > joined(short)


def test_inspect_compact_is_bounded_text(tmp_path):
    text = inspect_compact(_sample_pptx(tmp_path))
    assert text.startswith("deck: ")
    assert "[S1]" in text
    # スライドあたりの上限（見出し行＋図形行）が効いている。
    for block in text.split("\n\n")[1:]:
        assert len(block) <= 1600 + 200  # 見出し行と omitted 行の余裕分


def test_inspect_compact_caps_slides(tmp_path):
    # 40枚のデッキ → 30枚で打ち切り、省略数を明記する。
    dsl = "\n---\n".join(f'slide statement\n  headline "S{i}"' for i in range(40))
    out = tmp_path / "big.pptx"
    out.write_bytes(slidegen.render_to_bytes(dsl))
    text = inspect_compact(str(out))
    assert "[S30]" in text and "[S31]" not in text
    assert "(10 more slides omitted)" in text


# --- TABLE/CHART/GROUP/箇条書き抽出（slidegen 以外が作った "普通の" pptx を想定） -----

def _pptx_with_table_chart_group_bullets(tmp_path) -> str:
    """python-pptx で直接、表・ネイティブチャート・グループ・箇条書きを含む1枚を組み立てる
    （slidegen が生成した pptx ではなく、任意の既存デッキを想定した実験用フィクスチャ）。"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tbl_shape = slide.shapes.add_table(2, 2, Inches(0.5), Inches(0.5), Inches(3), Inches(1))
    tbl = tbl_shape.table
    tbl.cell(0, 0).text = "項目"
    tbl.cell(0, 1).text = "値"
    tbl.cell(1, 0).text = "売上"
    tbl.cell(1, 1).text = "120"

    cd = CategoryChartData()
    cd.categories = ["Q1", "Q2", "Q3"]
    cd.add_series("売上", (120, 150, 135))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(4), Inches(0.5), Inches(4), Inches(3), cd,
    )

    grp = slide.shapes.add_group_shape()
    box1 = grp.shapes.add_textbox(Inches(0), Inches(0), Inches(2), Inches(1))
    box1.text_frame.text = "グループ内テキストA"
    box2 = grp.shapes.add_textbox(Inches(3), Inches(0), Inches(2), Inches(1))
    box2.text_frame.text = "グループ内テキストB"
    grp.left, grp.top, grp.width, grp.height = Inches(1), Inches(4), Inches(6), Inches(2)

    bullets = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(4), Inches(1))
    tf = bullets.text_frame
    tf.text = "見出し"
    p2 = tf.add_paragraph()
    p2.text = "サブ項目1"
    p2.level = 1
    p3 = tf.add_paragraph()
    p3.text = "サブ項目2"
    p3.level = 1

    out = tmp_path / "external.pptx"
    prs.save(str(out))
    return str(out)


def test_inspect_extracts_table_cells(tmp_path):
    data = inspect(_pptx_with_table_chart_group_bullets(tmp_path))
    shapes = data["slides"][0]["shapes"]
    tables = [s for s in shapes if s.get("table")]
    assert len(tables) == 1
    assert tables[0]["table"] == [["項目", "値"], ["売上", "120"]]


def test_inspect_extracts_chart_data(tmp_path):
    data = inspect(_pptx_with_table_chart_group_bullets(tmp_path))
    shapes = data["slides"][0]["shapes"]
    charts = [s for s in shapes if s.get("chart")]
    assert len(charts) == 1
    c = charts[0]["chart"]
    assert c["chart_type"].startswith("COLUMN_CLUSTERED")
    assert c["categories"] == ["Q1", "Q2", "Q3"]
    assert c["series"]["売上"] == [120.0, 150.0, 135.0]


def test_inspect_flattens_group_with_absolute_coords(tmp_path):
    data = inspect(_pptx_with_table_chart_group_bullets(tmp_path))
    shapes = data["slides"][0]["shapes"]
    # group 自体は出力されない（中身だけが展開される）。
    assert not any(s["kind"].startswith("GROUP") for s in shapes)
    texts = {s["text"] for s in shapes if s["text"]}
    assert "グループ内テキストA" in texts
    assert "グループ内テキストB" in texts
    # グループ内の2つ目のテキストボックスは1つ目よりスライド上で右にある
    # （グループの絶対配置 + 内部オフセットが正しく合成されている＝chOff/chExt変換が効いている）。
    a = next(s for s in shapes if s["text"] == "グループ内テキストA")
    b = next(s for s in shapes if s["text"] == "グループ内テキストB")
    assert b["x%"] > a["x%"]
    # グループの絶対配置(top=4in)通りに子要素が収まっている（スライド高に対する% で判定）。
    prs = Presentation(_pptx_with_table_chart_group_bullets(tmp_path))
    expected_y_pct = round(100.0 * Inches(4) / prs.slide_height, 1)
    assert abs(a["y%"] - expected_y_pct) < 1.0


def test_inspect_preserves_bullet_levels(tmp_path):
    data = inspect(_pptx_with_table_chart_group_bullets(tmp_path))
    shapes = data["slides"][0]["shapes"]
    bullet_shape = next(s for s in shapes if "見出し" in s["text"])
    assert bullet_shape["text"] == "- 見出し / -- サブ項目1 / -- サブ項目2"


def test_inspect_compact_includes_table_and_chart_data(tmp_path):
    text = inspect_compact(_pptx_with_table_chart_group_bullets(tmp_path))
    assert "table=[項目 / 値] ; [売上 / 120]" in text
    assert "chart=COLUMN_CLUSTERED" in text
    assert "cats=[Q1 / Q2 / Q3]" in text
    assert "売上=[120.0 / 150.0 / 135.0]" in text
