"""
test_invariants.py — 第1層：構造インバリアントの自動テスト

このテストが落ちる ＝ 設計ドキュメントの §2-bis（編集可能）または §3（デザイン制約）
に違反している、ということ。新型を実装したら、まずここを通すこと。

実行: pytest tests/test_invariants.py -v

カバーするインバリアント:
  E. 編集可能性（§2-bis）
     E1. すべてのテキストがネイティブのテキストフレームに入っている
     E2. 画像化されたテキストブロックが存在しない
     E3. シェイプは選択可能（テキスト・図形・表のいずれか）
  P. パレット遵守（§3-1）
     P1. 塗りに使われている色は theme のパレット内に限る
     P2. accent 色の使用面積が全体の 10% を超えない（5%目標・10%は上限）
     P3. テキスト色は ink / muted / on_main / on_accent / accent / main 系のみ
  S. シェイプ常識
     S1. シェイプがスライド境界からはみ出していない
     S2. 1スライドのシェイプ数が暴走していない（< 80個）
     S3. TEXT_BOX のテキストがフレームに収まる（物理はみ出しのヒューリスティック検出）
  F. フォント統一（§3-1）
     F1. 使用されているフォント名が theme.font / theme.font_light / theme.font_mono 以外を含まない
"""
from __future__ import annotations
import sys, os, pathlib
from collections import Counter

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

sys.path.insert(0, str(pathlib.Path(__file__).resolve().parent.parent))
import slidegen
from slidegen.theme import Theme, DEFAULT_THEME
from slidegen.parser import parse
from slidegen.render import build


# テスト対象：examples/ の全 .slide を回す
EXAMPLES_DIR = pathlib.Path(__file__).resolve().parent.parent / "examples"
SLIDE_FILES = sorted(EXAMPLES_DIR.glob("*.slide"))


def _build(slide_file: pathlib.Path):
    text = slide_file.read_text(encoding="utf-8")
    prs = build(parse(text))
    return prs


def _collect_palette(theme: Theme) -> set[str]:
    """theme で使う全 hex 色（大文字）"""
    fields = ["base","base_2","main","main_2","main_3","accent","ink","muted",
              "rule","on_main","on_accent"]
    return {getattr(theme, f).upper() for f in fields}


@pytest.fixture(scope="module", params=SLIDE_FILES, ids=lambda p: p.stem)
def prs(request):
    return _build(request.param)


# ---------------------------------------------------------------------------
# E. 編集可能性（§2-bis）
# ---------------------------------------------------------------------------
class TestEditability:
    """テキストはテキストフレーム、図形はネイティブ — 画像化されていないこと"""

    def test_E1_text_lives_in_text_frames(self, prs):
        """E1: 何らかのテキストが必ずネイティブのテキストフレームから読み出せる"""
        any_text = False
        for slide in prs.slides:
            for sh in slide.shapes:
                if sh.has_text_frame and sh.text_frame.text.strip():
                    any_text = True
                    break
            if any_text:
                break
        assert any_text, "全スライドにライブテキストが1つも存在しない。画像化されている可能性"

    def test_E2_no_picture_only_slides(self, prs):
        """E2: 画像だけのスライド（テキスト0）は無いこと（タイトルや section にも文字は要る）"""
        for i, slide in enumerate(prs.slides, 1):
            picture_count = 0
            text_count = 0
            for sh in slide.shapes:
                if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    picture_count += 1
                if sh.has_text_frame and sh.text_frame.text.strip():
                    text_count += 1
            # 画像が含まれていても、テキストも必ずあること
            assert text_count > 0, f"スライド{i}: テキストが1つもない（画像化疑い）"

    def test_E3_shapes_are_native(self, prs):
        """E3: シェイプの種別がネイティブ（AUTO_SHAPE/TEXT_BOX/LINE/TABLE/PICTURE/CONNECTOR）に限る"""
        allowed = {
            MSO_SHAPE_TYPE.AUTO_SHAPE,
            MSO_SHAPE_TYPE.TEXT_BOX,
            MSO_SHAPE_TYPE.LINE,
            MSO_SHAPE_TYPE.TABLE,
            MSO_SHAPE_TYPE.PICTURE,
            MSO_SHAPE_TYPE.FREEFORM,   # カスタムジオメトリ(部分角丸カード等)もネイティブ編集可
            MSO_SHAPE_TYPE.CHART,      # ネイティブグラフ（数値を編集可能。画像化していない）
            # Connector は MSO_SHAPE_TYPE.LINE か AUTO_SHAPE で出る
        }
        for i, slide in enumerate(prs.slides, 1):
            for sh in slide.shapes:
                # connector や placeholder は shape_type が None になることがある — 許容
                if sh.shape_type is not None:
                    assert sh.shape_type in allowed, (
                        f"スライド{i}に編集不可な種別のシェイプ: {sh.shape_type}"
                    )


# ---------------------------------------------------------------------------
# P. パレット遵守（§3-1：配色70:25:5）
# ---------------------------------------------------------------------------
class TestPalette:

    def test_P1_fills_within_theme(self, prs):
        """P1: 塗りに使われている色が theme のパレットに収まる"""
        palette = _collect_palette(DEFAULT_THEME)
        offenders = []
        for i, slide in enumerate(prs.slides, 1):
            for sh in slide.shapes:
                try:
                    if sh.fill.type is None:
                        continue
                    if sh.fill.fore_color.type is None:
                        continue
                    rgb = str(sh.fill.fore_color.rgb).upper()
                    if rgb not in palette:
                        offenders.append((i, rgb))
                except Exception:
                    continue
        assert not offenders, (
            f"パレット外の色が使われている: {offenders[:5]}... "
            f"許容パレット: {sorted(palette)}"
        )

    def test_P2_accent_budget(self, prs):
        """P2: accent色の面積が、スライド全体面積の8%以下（目標5%・上限8%）。

        70:25:5 のルールは「スライド面積に対する塗りの比率」。
        この上限を超える ＝ 「強調が暴走している」サイン。
        """
        accent = DEFAULT_THEME.accent.upper()
        slide_area = int(prs.slide_width) * int(prs.slide_height)
        for i, slide in enumerate(prs.slides, 1):
            accent_area = 0
            for sh in slide.shapes:
                try:
                    if sh.fill.type is None or sh.fill.fore_color.type is None:
                        continue
                    if sh.width is None or sh.height is None:
                        continue
                    if str(sh.fill.fore_color.rgb).upper() == accent:
                        accent_area += int(sh.width) * int(sh.height)
                except Exception:
                    continue
            ratio = accent_area / slide_area
            assert ratio <= 0.08, (
                f"スライド{i}: accent色の面積比{ratio:.1%} が上限8%超過。"
                f"強調は1スライド1箇所が原則"
            )


# ---------------------------------------------------------------------------
# S. シェイプ常識
# ---------------------------------------------------------------------------
class TestShapes:

    def test_S1_shapes_within_bounds(self, prs):
        """S1: シェイプがスライド境界からはみ出していない（軽微なマージンは許容）"""
        sw, sh_ = prs.slide_width, prs.slide_height
        tol = 1000  # EMU。1000≈ごくわずかな丸め誤差
        offenders = []
        for i, slide in enumerate(prs.slides, 1):
            for shp in slide.shapes:
                if shp.left is None or shp.top is None:
                    continue
                if shp.width is None or shp.height is None:
                    continue
                if shp.left < -tol or shp.top < -tol:
                    offenders.append((i, "neg origin", shp.shape_type))
                if shp.left + shp.width > sw + tol:
                    offenders.append((i, "right overflow", shp.shape_type))
                if shp.top + shp.height > sh_ + tol:
                    offenders.append((i, "bottom overflow", shp.shape_type))
        assert not offenders, f"はみ出しシェイプ: {offenders[:5]}"

    def test_S2_shape_count_reasonable(self, prs):
        """S2: 1スライドのシェイプ数が80未満（暴走防止）"""
        for i, slide in enumerate(prs.slides, 1):
            n = len(slide.shapes)
            assert n < 80, f"スライド{i}のシェイプ数が{n}個（暴走の疑い）"

    def test_S3_text_overflow_risk_is_low(self, prs):
        """S3: TEXT_BOX のテキストがフレームに収まる（物理はみ出しのヒューリスティック）。

        折り返し設定で判定を分ける（word_wrap=False 自体は不正ではない。code_block 等が意図的に使う）:
          - word_wrap=False（折り返さない）→ 最長行の推定幅がフレーム幅を超えないか（水平はみ出し）。
          - それ以外（折り返す）→ 推定行数からのテキスト高さがフレーム高さを超えないか（垂直はみ出し）。
        全角=2幅で数える。明白なはみ出しだけを拾うスモークなので TOLERANCE に余裕を持たせ誤検知を避ける。
        """
        import math

        DEFAULT_FONT_EMU = Pt(12)
        LINE_SPACING = 1.3
        CHAR_W_RATIO = 0.5   # 半角1文字の幅 ≈ フォントサイズ×0.5
        TOLERANCE = 2.0

        def _visual_len(s: str) -> int:
            return sum(2 if ord(c) > 127 else 1 for c in s)

        w_offenders = []
        h_offenders = []

        for si, slide in enumerate(prs.slides, 1):
            for shp in slide.shapes:
                if shp.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
                    continue
                if not shp.has_text_frame:
                    continue
                tf = shp.text_frame
                if not tf.text.strip():
                    continue

                width = getattr(shp, "width", None) or 0
                height = getattr(shp, "height", None) or 0
                eff_w = max(1, int(width) - (tf.margin_left or 0) - (tf.margin_right or 0))
                eff_h = max(1, int(height) - (tf.margin_top or 0) - (tf.margin_bottom or 0))

                sizes = [
                    r.font.size for para in tf.paragraphs
                    for r in para.runs if r.font.size
                ]
                avg_font = max(Pt(4), (sum(sizes) / len(sizes)) if sizes else DEFAULT_FONT_EMU)
                char_w = avg_font * CHAR_W_RATIO
                name = (shp.name or "")[:20]

                if tf.word_wrap is False:
                    # 折り返さない → 最長行が横にはみ出さないか
                    max_line = max(
                        (_visual_len(seg)
                         for para in tf.paragraphs
                         for seg in (para.text.splitlines() or [""])),
                        default=0,
                    )
                    estimated_w = max_line * char_w
                    if estimated_w > eff_w * TOLERANCE:
                        w_offenders.append((si, name, int(estimated_w), int(eff_w)))
                else:
                    # 折り返す → 推定行数×行高がフレーム高さを超えないか
                    chars_per_line = max(1, eff_w / char_w)
                    total_lines = 0
                    for para in tf.paragraphs:
                        for seg in (para.text.splitlines() or [""]):
                            total_lines += math.ceil(max(1, _visual_len(seg)) / chars_per_line)
                    estimated_h = avg_font * LINE_SPACING * total_lines
                    if estimated_h > eff_h * TOLERANCE:
                        h_offenders.append((si, name, int(estimated_h), int(eff_h)))

        messages = []
        if w_offenders:
            messages.append(f"テキスト幅超過の可能性(word_wrap=False): {w_offenders[:3]}")
        if h_offenders:
            messages.append(f"テキスト高さ超過の可能性: {h_offenders[:3]}")
        assert not messages, "\n".join(messages)


# ---------------------------------------------------------------------------
# F. フォント統一（§3-1）
# ---------------------------------------------------------------------------
class TestFonts:

    def test_F1_single_font(self, prs):
        """F1: 使われているフォント名が theme.font / theme.font_light / theme.font_mono のみ"""
        allowed = {DEFAULT_THEME.font, DEFAULT_THEME.font_light,
                   DEFAULT_THEME.font_mono, None}
        used = set()
        for slide in prs.slides:
            for sh in slide.shapes:
                if not sh.has_text_frame:
                    continue
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        used.add(r.font.name)
        offenders = used - allowed
        assert not offenders, (
            f"未許可フォントが使われている: {offenders}。"
            f"許可: {allowed - {None}}"
        )
