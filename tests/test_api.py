"""
test_api.py — バックエンド用 public API（slidegen/api.py）のテスト。

このAPIは「アプリのバックエンドから import して使う」入口。ディスクを介さず
メモリで pptx の bytes を得られること（render_to_bytes）が肝。
"""
import io
import pathlib

import pytest
from pptx import Presentation

REPO_ROOT = pathlib.Path(__file__).resolve().parent.parent
SAMPLE = (REPO_ROOT / "examples" / "sample.slide").read_text(encoding="utf-8")


def test_render_text_returns_presentation_with_expected_slide_count():
    import slidegen
    from slidegen.parser import parse

    prs = slidegen.render_text(SAMPLE)
    assert isinstance(prs, Presentation().__class__)
    assert len(prs.slides) == len(parse(SAMPLE))


def test_render_to_bytes_is_nonempty_zip():
    import slidegen

    data = slidegen.render_to_bytes(SAMPLE)
    assert isinstance(data, (bytes, bytearray))
    assert len(data) > 0
    # pptx は zip コンテナ。先頭は ZIP のローカルファイルシグネチャ。
    assert data[:4] == b"PK\x03\x04"


def test_render_to_bytes_roundtrips_back_into_presentation():
    import slidegen

    data = slidegen.render_to_bytes(SAMPLE)
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) > 0
    # 少なくとも1スライドに編集可能なシェイプが載っている。
    assert any(len(s.shapes) > 0 for s in prs.slides)


def test_render_file_writes_pptx(tmp_path):
    import slidegen

    out = tmp_path / "deck.pptx"
    src = tmp_path / "deck.slide"
    src.write_text(SAMPLE, encoding="utf-8")

    returned = slidegen.render_file(src, out)
    assert pathlib.Path(returned) == out
    assert out.exists() and out.stat().st_size > 0
    Presentation(str(out))  # 開けること


def test_api_works_without_top_level_import():
    """`import slidegen.api` のように親パッケージのトップ import を介さない経路でも、
    親 __init__ が先に走って全 render_* が登録されるため、全型がレンダリングできる。"""
    import importlib

    api = importlib.import_module("slidegen.api")
    # examples/showcase.slide は多様な型を含む。型登録が抜けていれば未対応型で落ちる。
    showcase = (REPO_ROOT / "examples" / "showcase.slide").read_text(encoding="utf-8")
    data = api.render_to_bytes(showcase)
    assert data[:4] == b"PK\x03\x04"
