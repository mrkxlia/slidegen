"""
render_tech.py — 技術資料の個別型。

- code_block   : コードを等幅フォントで表示（濃色パネル＋行）。シンタックスハイライトはしない（崩れ防止）。
- terminal     : ターミナル出力風（黒地＋プロンプト）。
- api_endpoint_table : APIエンドポイント一覧（メソッド色分け＋パス＋説明）。
- code_diff    : 差分表示（行頭 +/- で文字色を変える。S5c）。
- sql_result   : クエリ（等幅パネル）＋結果テーブル（本物のPowerPointテーブル。S5c）。

設計思想：標準図形のみ。等幅は theme.font_mono。色は theme 経由。
"""
from __future__ import annotations
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from . import render as R
from .render import (add_rect, add_text, render_header, render_foot,
                     SLIDE_W, SLIDE_H, MARGIN, CONTENT_W)
from .parser import Slide


def _code_lines(data: Slide):
    """本文のコード行を取得。最初のブロックの lines を使う。"""
    if data.blocks:
        return data.blocks[0].lines
    return []


def _mono_text(slide, x, y, w, h, theme, lines, *, size, color_name, on_dark, line_colors=None):
    """等幅テキストを1行=1段落で描く。line_colors を渡すと行ごとに色を変えられる
    （未指定時は全行 color_name の単色＝既存呼び出し元の挙動は不変）。"""
    box = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(6)
    tf.margin_top = tf.margin_bottom = Pt(4)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for j, ln in enumerate(lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = ln if ln else " "
        r.font.size = Pt(size)
        r.font.name = theme.font_mono
        c = line_colors[j] if line_colors and j < len(line_colors) else color_name
        r.font.color.rgb = theme.rgb(c)
    return box


def render_code_block(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    lines = _code_lines(data)
    bottom = SLIDE_H - Inches(0.9)
    y = top + Inches(0.1)
    h = bottom - y
    # 濃色パネル
    add_rect(slide, MARGIN, int(y), CONTENT_W, int(h), theme, "ink", rounded=True)
    # ファイル名/言語（任意）
    lang = data.props.get("lang") or data.props.get("file")
    pad = Inches(0.2)
    ty = y + pad
    if lang:
        add_text(slide, MARGIN + pad, int(ty), CONTENT_W - pad * 2, Inches(0.3), theme,
                 lang, size=11, color_name="muted", bold=True)
        ty += Inches(0.35)
    _mono_text(slide, MARGIN + pad, ty, CONTENT_W - pad * 2, bottom - ty - Inches(0.1),
               theme, lines, size=14, color_name="base", on_dark=True)


def render_terminal(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    lines = _code_lines(data)
    bottom = SLIDE_H - Inches(0.9)
    y = top + Inches(0.1)
    h = bottom - y
    add_rect(slide, MARGIN, int(y), CONTENT_W, int(h), theme, "ink", rounded=True)
    # 行頭に $ を付ける（コマンド行の見立て）。出力行は先頭にスペースを入れて区別する運用。
    shown = []
    for ln in lines:
        if ln.startswith("  ") or ln.startswith("\t"):
            shown.append(ln)            # 出力行はそのまま
        else:
            shown.append("$ " + ln)     # コマンド行
    pad = Inches(0.2)
    _mono_text(slide, MARGIN + pad, y + pad, CONTENT_W - pad * 2, h - pad * 2,
               theme, shown, size=14, color_name="base", on_dark=True)


# メソッド → 色
_METHOD_COLOR = {"GET": "main_2", "POST": "main", "PUT": "muted",
                 "PATCH": "muted", "DELETE": "accent"}


def render_api_endpoint_table(slide, data: Slide, theme):
    """各 col が1エンドポイント。
      col "GET"            # title = HTTPメソッド
        "/api/users"       # line[0] = パス
        "ユーザー一覧取得"  # line[1] = 説明
    """
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    rows = data.blocks
    n = len(rows)
    if n == 0:
        return
    bottom = SLIDE_H - Inches(0.8)
    y0 = top + Inches(0.2)
    rh = min(Inches(0.7), (bottom - y0 - Inches(0.1) * n) / max(n, 1))
    gap = Inches(0.1)
    method_w = Inches(1.3)
    path_w = Inches(4.2)

    for i, b in enumerate(rows):
        y = y0 + i * (rh + gap)
        method = (b.title or "GET").strip().upper()
        color = _METHOD_COLOR.get(method, "main")
        # メソッドバッジ
        add_rect(slide, MARGIN, int(y), int(method_w), int(rh), theme, color, rounded=True)
        add_text(slide, MARGIN, int(y), int(method_w), int(rh), theme, method,
                 size=13, color_name="on_main", bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # パス（等幅）
        path = b.lines[0] if b.lines else ""
        pbox = slide.shapes.add_textbox(int(MARGIN + method_w + gap), int(y),
                                        int(path_w), int(rh))
        tf = pbox.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pr = tf.paragraphs[0].add_run(); pr.text = path
        pr.font.name = theme.font_mono; pr.font.size = Pt(14); pr.font.color.rgb = theme.rgb("ink")
        # 説明
        desc = b.lines[1] if len(b.lines) > 1 else ""
        dx = MARGIN + method_w + gap + path_w + gap
        add_text(slide, int(dx), int(y), int(SLIDE_W - dx - MARGIN), int(rh), theme,
                 desc, size=13, color_name="muted",
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------------------
# code_diff — 差分表示（行頭 +/- で文字色を変える。プレフィックスは表示に残す）
#   col の各行の先頭文字で判定： "+"=追加(main_2) / "-"=削除(accent) / それ以外=文脈(base)
# ---------------------------------------------------------------------------
def render_code_diff(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    lines = _code_lines(data)
    bottom = SLIDE_H - Inches(0.9)
    y = top + Inches(0.1)
    h = bottom - y
    add_rect(slide, MARGIN, int(y), CONTENT_W, int(h), theme, "ink", rounded=True)
    lang = data.props.get("lang") or data.props.get("file")
    pad = Inches(0.2)
    ty = y + pad
    if lang:
        add_text(slide, MARGIN + pad, int(ty), CONTENT_W - pad * 2, Inches(0.3), theme,
                 lang, size=11, color_name="muted", bold=True)
        ty += Inches(0.35)
    colors = []
    for ln in lines:
        if ln.startswith("+"):
            colors.append("main_2")
        elif ln.startswith("-"):
            colors.append("accent")
        else:
            colors.append("base")
    _mono_text(slide, MARGIN + pad, ty, CONTENT_W - pad * 2, bottom - ty - Inches(0.1),
               theme, lines, size=14, color_name="base", on_dark=True, line_colors=colors)


# ---------------------------------------------------------------------------
# sql_result — クエリ（等幅パネル）＋結果テーブル（本物のPowerPointテーブル）
#   query "SELECT ..." "FROM ..." ... （複数値propで複数行のクエリを受ける）
#   col "列名"    # 結果テーブルの1列。lines=その列の値（上から順に行）
# ---------------------------------------------------------------------------
_SQL_MAX_COLS = 6
_SQL_MAX_ROWS = 8


def render_sql_result(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    bottom = SLIDE_H - Inches(0.8)

    query_lines = data.props.get("query_list")
    if query_lines is None:
        single = data.props.get("query")
        query_lines = [single] if single else []

    y = top + Inches(0.1)
    if query_lines:
        panel_h = min(Inches(0.35) * len(query_lines) + Inches(0.2), Inches(1.8))
        add_rect(slide, MARGIN, int(y), CONTENT_W, int(panel_h), theme, "ink", rounded=True)
        _mono_text(slide, MARGIN + Inches(0.2), y + Inches(0.1), CONTENT_W - Inches(0.4),
                   panel_h - Inches(0.2), theme, query_lines, size=13, color_name="base", on_dark=True)
        y += panel_h + Inches(0.2)

    cols = data.blocks[:_SQL_MAX_COLS]
    if not cols:
        return
    ncol = len(cols)
    nrow = 1 + min(max((len(c.lines) for c in cols), default=0), _SQL_MAX_ROWS)
    tbl_h = bottom - y
    gfx = slide.shapes.add_table(nrow, ncol, MARGIN, int(y), CONTENT_W, int(tbl_h))
    table = gfx.table
    for ci, blk in enumerate(cols):
        header_cell = table.cell(0, ci)
        header_cell.text = blk.title
        hp = header_cell.text_frame.paragraphs[0]
        hp.alignment = PP_ALIGN.CENTER
        for run in hp.runs:
            run.font.name = theme.font_mono
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = theme.rgb("on_main")
        header_cell.fill.solid()
        header_cell.fill.fore_color.rgb = theme.rgb("main")
        for ri in range(1, nrow):
            cell = table.cell(ri, ci)
            val = blk.lines[ri - 1] if ri - 1 < len(blk.lines) else ""
            cell.text = val
            cp = cell.text_frame.paragraphs[0]
            cp.alignment = PP_ALIGN.CENTER
            for run in cp.runs:
                run.font.name = theme.font_mono
                run.font.size = Pt(12)
                run.font.color.rgb = theme.rgb("ink")
            cell.fill.solid()
            cell.fill.fore_color.rgb = theme.rgb("base") if ri % 2 else theme.rgb("base_2")


R.register("code_block", render_code_block)
R.register("terminal", render_terminal)
R.register("api_endpoint_table", render_api_endpoint_table)
R.register("code_diff", render_code_diff)
R.register("sql_result", render_sql_result)
