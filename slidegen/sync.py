"""
sync.py — 手編集同期（Hand-Edit Sync）。

流れ：
  1) slidegen が .slide から pptx を生成（これが「ベースライン」）
  2) 人が PowerPoint で pptx のテキストを修正
  3) この sync が「ベースラインのテキスト」と「修正後 pptx のテキスト」を
     スライド単位で突き合わせ、変わった箇所を検出
  4) 変わったテキストで元の .slide を書き換える（テキストの差し替えのみ）

設計方針（割り切り）：
  - 双方向の完全逆変換はやらない。壊れやすいため。
  - 「テキストの差し替え」に限定する。図形の移動・追加・色変更は対象外（記法に存在しないため）。
  - ベースラインは .slide を再レンダリングして得る（生成時の文字列が分かる）。
  - マッチングは「スライドindex × 出現順 × 文字列の近さ」で行う。

これにより、Sonnet生成 → 人が文言を直す → 直しを記法に取り込む、のループが回る。

使い方:
  python -m slidegen.sync original.slide edited.pptx            # 差分を表示（dry-run）
  python -m slidegen.sync original.slide edited.pptx --apply    # .slide を書き換える
  python -m slidegen.sync original.slide edited.pptx --apply -o new.slide  # 別名で保存
"""
from __future__ import annotations
import argparse
import difflib
from pptx import Presentation

from .parser import parse
from .render import build


def _shape_sort_key(sh):
    """シェイプを上→下、左→右の順で安定ソートするキー。"""
    top = sh.top if sh.top is not None else 0
    left = sh.left if sh.left is not None else 0
    # 行が近いものは同じ行とみなして左→右（Emu約30万=0.3inchで丸め）
    band = int(top) // 300000
    return (band, int(left))


def extract_slide_texts(prs) -> list[list[str]]:
    """各スライドの非空テキストを、出現順（上→下、左→右）に並べて返す。
    戻り値: [ [slide1のテキスト...], [slide2のテキスト...], ... ]
    """
    out = []
    for slide in prs.slides:
        texts = []
        shapes = sorted(
            [sh for sh in slide.shapes if sh.has_text_frame],
            key=_shape_sort_key,
        )
        for sh in shapes:
            t = sh.text_frame.text
            if t and t.strip():
                texts.append(t.strip())
        out.append(texts)
    return out


def _texts_from_slide_source(slide_source: str) -> list[list[str]]:
    """.slide を再レンダリングして、生成時のテキスト（ベースライン）を得る。"""
    slides = parse(slide_source)
    prs = build(slides)
    return extract_slide_texts(prs)


def _texts_from_pptx_path(path: str) -> list[list[str]]:
    prs = Presentation(path)
    return extract_slide_texts(prs)


def _match_changes(baseline: list[str], edited: list[str]):
    """1スライド内で baseline → edited の変更ペアを返す。
    戻り値: [(old_text, new_text), ...]（変わったものだけ）
    difflib で行単位の対応を取り、置換だけを拾う。
    """
    changes = []
    sm = difflib.SequenceMatcher(a=baseline, b=edited, autojunk=False)
    for tag, i1, i2, j1, j2 in sm.get_opcodes():
        if tag == "replace":
            # 同数の置換は1対1で対応づける
            old_block = baseline[i1:i2]
            new_block = edited[j1:j2]
            for k in range(min(len(old_block), len(new_block))):
                if old_block[k] != new_block[k]:
                    changes.append((old_block[k], new_block[k]))
            # 余り（数が違う場合）は安全のため触らない（報告のみ別途）
    return changes


def compute_diff(slide_source: str, edited_pptx_path: str):
    """ベースラインと編集後を突き合わせ、スライドごとの変更リストを返す。
    戻り値: [ {slide: i, changes: [(old,new),...], note: str}, ... ]
    """
    base = _texts_from_slide_source(slide_source)
    edited = _texts_from_pptx_path(edited_pptx_path)
    results = []
    n = max(len(base), len(edited))
    for i in range(n):
        b = base[i] if i < len(base) else []
        e = edited[i] if i < len(edited) else []
        changes = _match_changes(b, e)
        note = ""
        if len(b) != len(e):
            note = f"テキスト要素数が変化（生成{len(b)}→編集{len(e)}）。増減分は自動反映しません。"
        if changes or note:
            results.append({"slide": i + 1, "changes": changes, "note": note})
    return results


def apply_changes(slide_source: str, diff) -> tuple[str, int]:
    """差分を .slide ソース文字列に適用（テキストの差し替えのみ）。
    戻り値: (新しいsource, 適用件数)
    """
    new_src = slide_source
    applied = 0
    for entry in diff:
        for old, new in entry["changes"]:
            # .slide 内に old がそのまま含まれていれば置換（最初の1件）
            if old and old in new_src and old != new:
                new_src = new_src.replace(old, new, 1)
                applied += 1
    return new_src, applied


def _format_diff(diff) -> str:
    if not diff:
        return "変更は検出されませんでした（テキストは生成時と同一）。"
    lines = []
    for entry in diff:
        lines.append(f"■ Slide {entry['slide']}")
        for old, new in entry["changes"]:
            lines.append(f"   - 変更前: {old}")
            lines.append(f"   + 変更後: {new}")
        if entry["note"]:
            lines.append(f"   ⚠ {entry['note']}")
    return "\n".join(lines)


def main():
    ap = argparse.ArgumentParser(
        description="手編集同期：編集後pptxの文言変更を元の.slideに反映")
    ap.add_argument("slide", help="生成元の記法ファイル(.slide)")
    ap.add_argument("edited_pptx", help="人が編集した後の .pptx")
    ap.add_argument("--apply", action="store_true", help=".slide を実際に書き換える")
    ap.add_argument("-o", "--output", default=None,
                    help="書き換え結果の保存先（省略時は元の.slideを上書き）")
    args = ap.parse_args()

    with open(args.slide, encoding="utf-8") as f:
        src = f.read()

    diff = compute_diff(src, args.edited_pptx)
    print(_format_diff(diff))

    if args.apply:
        new_src, applied = apply_changes(src, diff)
        out = args.output or args.slide
        with open(out, "w", encoding="utf-8") as f:
            f.write(new_src)
        print(f"\n→ {applied}件の文言変更を {out} に反映しました。")
    else:
        total = sum(len(e["changes"]) for e in diff)
        if total:
            print(f"\n（dry-run）{total}件の文言変更を検出。--apply で .slide に反映します。")


if __name__ == "__main__":
    main()
