"""
render_tech_diagrams.py — 技術資料の図解系個別型（S5c）。

- layered_stack   : 技術スタック（層を上から積む。等幅版pyramid）
- c4_context      : C4コンテキスト図（中心システム1＋周辺アクター/外部システムN。tree の
                    ジオメトリを踏襲するが中心は常にaccent＝C4の意味論に合わせた独立実装）
- sequence_diagram: シーケンス図（参加者＝縦のライフライン、メッセージ＝横向きブロック矢印）
- state_transition: 状態遷移図（状態を円周上に配置、遷移は直線＋ラベル。矢頭は使わない）
- er_diagram      : ER図（エンティティ＝属性リストカード、リレーション＝直線＋カーディナリティ）

設計思想：標準プリセット図形のみ・回転禁止（type_catalog.md §6）。sequence_diagram/
state_transition/er_diagram も Mermaid でレンダリングして画像を貼ることはしない
（ADR 0002 の画像化絶対禁止と正面衝突するため）。標準図形（矩形・直線コネクタ・テキスト）の
組み合わせのみで表現する。

sequence_diagram / state_transition / er_diagram に共通する記法規約：
  「from/to の rows を持つ col ＝ 接続（メッセージ/遷移/リレーション）を表すブロック」。
  er_diagram だけは「from/to を持たない col」もあり、それはエンティティ（属性リスト）として
  扱う。
"""
from __future__ import annotations
import math
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from . import render as R
from .render import (add_rect, add_text, render_header, render_foot,
                     SLIDE_W, SLIDE_H, MARGIN, CONTENT_W)
from .parser import Slide
from .render_util import block_items, add_items_text, columns_geometry, fill_shape


def _line(slide, theme, x1, y1, x2, y2, color="rule", weight=1.5):
    ln = slide.shapes.add_connector(2, int(x1), int(y1), int(x2), int(y2))
    ln.line.color.rgb = theme.rgb(color)
    ln.line.width = Pt(weight)
    return ln


def _row_dict(blk):
    """block.rows を {label: value} にする（from/to 読み出し用）。"""
    return dict(blk.rows)


def _rect_edge_point(rect, tx, ty):
    """rect=(x,y,w,h) の中心から (tx,ty) 方向へ伸ばした半直線が rect の境界と交わる点。
    中心-中心で線を引くとカードに隠れてしまう(er_diagram等)ため、線はこの境界点同士を結ぶ。"""
    x, y, w, h = rect
    cx, cy = x + w / 2, y + h / 2
    dx, dy = tx - cx, ty - cy
    if dx == 0 and dy == 0:
        return cx, cy
    scale = float("inf")
    if dx != 0:
        scale = min(scale, abs((w / 2) / dx))
    if dy != 0:
        scale = min(scale, abs((h / 2) / dy))
    return cx + dx * scale, cy + dy * scale


def _add_outline_rect(slide, theme, x, y, w, h, color_name, weight=2.5):
    """塗りなし・枠線のみの矩形（render_charts_shapes.py の marimekko/treemap と同じ手法。
    線なので P2 accent算入対象外。面積の大きいセルを accent 塗りにするとP2上限を超えやすい
    layered_stack の highlight で使う）。"""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), int(w), int(h))
    shp.fill.background()
    shp.line.color.rgb = theme.rgb(color_name)
    shp.line.width = Pt(weight)
    shp.shadow.inherit = False
    return shp


# ---------------------------------------------------------------------------
# layered_stack — 技術スタック（col 記述順=上から。等幅段のpyramid）
# ---------------------------------------------------------------------------
_LAYERED_STACK_MAX = 6


def render_layered_stack(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    layers = data.blocks[:_LAYERED_STACK_MAX]
    n = len(layers)
    if n == 0:
        return
    bottom = SLIDE_H - Inches(0.7)
    avail_h = bottom - top - Inches(0.1)
    gap = Inches(0.08)
    layer_h = (avail_h - gap * (n - 1)) / n

    for i, blk in enumerate(layers):
        y = top + Inches(0.1) + i * (layer_h + gap)
        # 層は面積が大きく accent 塗りだとP2(8%上限)を超過しやすいため、highlight は
        # アウトライン枠線のみで表現する（塗りは常にmain。marimekko/treemapと同じ判断）。
        add_rect(slide, MARGIN, int(y), CONTENT_W, int(layer_h), theme, "main", rounded=True)
        if blk.highlight:
            _add_outline_rect(slide, theme, MARGIN, y, CONTENT_W, layer_h, "accent")
        add_text(slide, MARGIN + Inches(0.25), int(y), CONTENT_W * 0.35, int(layer_h), theme,
                 blk.title, size=16, color_name="on_main", bold=True,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        items = block_items(blk)
        if items:
            add_text(slide, MARGIN + CONTENT_W * 0.4, int(y), CONTENT_W * 0.55, int(layer_h),
                     theme, "  ".join(items), size=12, color_name="on_main",
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------------------
# c4_context — 中心システム1(常にaccent)＋周辺N。tree と同じ辺中央直線コネクタ。
# ---------------------------------------------------------------------------
_C4_CONTEXT_MAX_CHILDREN = 6


def render_c4_context(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    if not data.blocks:
        return
    bottom = SLIDE_H - Inches(0.7)
    h_total = bottom - top
    center = data.blocks[0]
    children = data.blocks[1:1 + _C4_CONTEXT_MAX_CHILDREN]

    node_w = Inches(2.6)
    node_h = Inches(0.9)
    px = SLIDE_W / 2 - node_w / 2
    py = top + Inches(0.2)

    if children:
        n = len(children)
        gap = Inches(0.3)
        cw = columns_geometry(CONTENT_W, n, gap)
        cy = py + node_h + Inches(0.9)
        for i, ch in enumerate(children):
            cx = MARGIN + i * (cw + gap)
            _line(slide, theme, px + node_w / 2, py + node_h, cx + cw / 2, cy)
            color = "main_2" if not ch.highlight else "accent"
            add_rect(slide, cx, cy, cw, node_h, theme, color, rounded=True)
            add_text(slide, cx, cy, cw, node_h, theme, ch.title,
                     size=13, color_name="on_main", bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            if ch.lines:
                add_text(slide, cx, cy + node_h + Inches(0.12), cw,
                         h_total - (cy + node_h + Inches(0.12) - top),
                         theme, " ".join(ch.lines), size=10, color_name="muted",
                         align=PP_ALIGN.CENTER)

    # 中心システムは常に強調（C4の意味論：中心=このコンテキストの主役）
    add_rect(slide, px, py, node_w, node_h, theme, "accent", rounded=True)
    add_text(slide, px, py, node_w, node_h, theme, center.title,
             size=16, color_name="on_accent", bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if center.lines:
        add_text(slide, px - Inches(0.4), py + node_h + Inches(0.02), node_w + Inches(0.8),
                 Inches(0.35), theme, " ".join(center.lines), size=10, color_name="muted",
                 align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# sequence_diagram — 参加者=縦ライフライン、メッセージ=横向きブロック矢印
#   participants "A" "B" "C"（上限5）
#   col   from "A"  to "B"  "メッセージ文"（上限10。記述順=時系列）
# ---------------------------------------------------------------------------
_SEQ_MAX_PARTICIPANTS = 5
_SEQ_MAX_MESSAGES = 10


def render_sequence_diagram(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    participants = data.props.get("participants_list")
    if participants is None:
        single = data.props.get("participants")
        participants = [single] if single else []
    participants = participants[:_SEQ_MAX_PARTICIPANTS]
    if len(participants) < 2:
        return
    messages = data.blocks[:_SEQ_MAX_MESSAGES]

    bottom = SLIDE_H - Inches(0.7)
    gap = Inches(0.2)
    head_h = Inches(0.5)
    pw = columns_geometry(CONTENT_W, len(participants), gap)
    centers = [MARGIN + i * (pw + gap) + pw / 2 for i in range(len(participants))]

    lifeline_bottom = bottom
    for i, name in enumerate(participants):
        x = MARGIN + i * (pw + gap)
        add_rect(slide, x, top, pw, head_h, theme, "main", rounded=True)
        add_text(slide, x, top, pw, head_h, theme, name,
                 size=13, color_name="on_main", bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _line(slide, theme, centers[i], top + head_h, centers[i], lifeline_bottom, "rule", 1.0)

    if not messages:
        return
    msg_top = top + head_h + Inches(0.25)
    row_h = (lifeline_bottom - msg_top - Inches(0.1)) / len(messages)
    name_to_idx = {name: i for i, name in enumerate(participants)}

    for i, blk in enumerate(messages):
        rows = _row_dict(blk)
        frm, to = rows.get("from"), rows.get("to")
        if frm not in name_to_idx or to not in name_to_idx or frm == to:
            continue  # 不明な参加者・自己メッセージ(v1非対応)はスキップ
        y = msg_top + i * row_h + row_h / 2
        fx, tx = centers[name_to_idx[frm]], centers[name_to_idx[to]]
        label = blk.lines[0] if blk.lines else ""
        color = "accent" if blk.highlight else "main_2"
        arrow_h = Inches(0.14)
        if label:
            add_text(slide, min(fx, tx), y - Inches(0.32), abs(tx - fx) or Inches(0.5),
                     Inches(0.28), theme, label, size=10, color_name="muted",
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
        shape = MSO_SHAPE.RIGHT_ARROW if tx >= fx else MSO_SHAPE.LEFT_ARROW
        w = abs(tx - fx)
        if w <= 0:
            continue
        shp = slide.shapes.add_shape(shape, int(min(fx, tx)), int(y - arrow_h / 2), int(w), int(arrow_h))
        fill_shape(shp, theme, color, no_shadow=True)


# ---------------------------------------------------------------------------
# state_transition — 状態を円周上に配置。遷移は直線＋ラベル（矢頭なし）。
#   states "A" "B" "C"（上限6）
#   col   from "A"  to "B"  "ラベル"（上限10）
# ---------------------------------------------------------------------------
_STATE_MAX_STATES = 6
_STATE_MAX_TRANSITIONS = 10


def render_state_transition(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    states = data.props.get("states_list")
    if states is None:
        single = data.props.get("states")
        states = [single] if single else []
    states = states[:_STATE_MAX_STATES]
    n = len(states)
    if n == 0:
        return
    transitions = data.blocks[:_STATE_MAX_TRANSITIONS]

    bottom = SLIDE_H - Inches(0.7)
    cx = SLIDE_W / 2
    cy = (top + bottom) / 2
    radius = min((bottom - top), CONTENT_W) / 2 - Inches(0.9)
    node_w, node_h = Inches(1.8), Inches(0.7)

    positions = {}
    for i, name in enumerate(states):
        angle = -math.pi / 2 + 2 * math.pi * i / n
        positions[name] = (cx + radius * math.cos(angle), cy + radius * math.sin(angle))

    name_highlight = {}
    node_rect = {name: (px - node_w / 2, py - node_h / 2, node_w, node_h)
                 for name, (px, py) in positions.items()}
    # 遷移線はボックス境界の点同士を結ぶ（中心-中心だとボックスの下に隠れて見えなくなる）
    for blk in transitions:
        rows = _row_dict(blk)
        frm, to = rows.get("from"), rows.get("to")
        if frm not in positions or to not in positions or frm == to:
            continue  # 未知の状態名・自己遷移(v1非対応)はスキップ
        fcx, fcy = positions[frm]
        tcx, tcy = positions[to]
        fx, fy = _rect_edge_point(node_rect[frm], tcx, tcy)
        tx, ty = _rect_edge_point(node_rect[to], fcx, fcy)
        color = "accent" if blk.highlight else "rule"
        _line(slide, theme, fx, fy, tx, ty, color, 1.5)
        label = blk.lines[0] if blk.lines else ""
        if label:
            mx, my = (fx + tx) / 2, (fy + ty) / 2
            add_text(slide, mx - Inches(0.8), my - Inches(0.15), Inches(1.6), Inches(0.3),
                     theme, label, size=10, color_name="muted",
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if blk.highlight:
            name_highlight[to] = True

    for name in states:
        x, y = positions[name]
        color = "accent" if name_highlight.get(name) else "main"
        add_rect(slide, x - node_w / 2, y - node_h / 2, node_w, node_h, theme, color, rounded=True)
        add_text(slide, x - node_w / 2, y - node_h / 2, node_w, node_h, theme, name,
                 size=13, color_name="on_main", bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------------------
# er_diagram — from/to を持たない col=エンティティ（属性リストカード）、
#              from/to を持つ col=リレーション（直線＋カーディナリティ）
# ---------------------------------------------------------------------------
_ER_MAX_ENTITIES = 8
_ER_MAX_RELATIONS = 10
_ER_COLS = 4


def render_er_diagram(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    entity_blocks, relation_blocks = [], []
    for blk in data.blocks:
        rows = _row_dict(blk)
        if "from" in rows and "to" in rows:
            relation_blocks.append(blk)
        else:
            entity_blocks.append(blk)
    entity_blocks = entity_blocks[:_ER_MAX_ENTITIES]
    relation_blocks = relation_blocks[:_ER_MAX_RELATIONS]
    n = len(entity_blocks)
    if n == 0:
        return

    bottom = SLIDE_H - Inches(0.7)
    avail_h = bottom - top
    gap = Inches(0.45)
    cols = min(_ER_COLS, n)
    erows = (n + cols - 1) // cols
    cw = columns_geometry(CONTENT_W, cols, gap)
    ch = (avail_h - gap * (erows - 1)) / erows if erows else avail_h

    rects = {}
    for i, blk in enumerate(entity_blocks):
        r, c = divmod(i, cols)
        x = MARGIN + c * (cw + gap)
        y = top + r * (ch + gap)
        rects[blk.title] = (x, y, cw, ch)

    # リレーション線はカード境界の点同士を結ぶ（中心-中心だとカードの下に隠れて見えなくなる）
    for blk in relation_blocks:
        rows = _row_dict(blk)
        frm, to = rows.get("from"), rows.get("to")
        if frm not in rects or to not in rects:
            continue
        fr, tr = rects[frm], rects[to]
        fcx, fcy = fr[0] + fr[2] / 2, fr[1] + fr[3] / 2
        tcx, tcy = tr[0] + tr[2] / 2, tr[1] + tr[3] / 2
        fx, fy = _rect_edge_point(fr, tcx, tcy)
        tx, ty = _rect_edge_point(tr, fcx, fcy)
        _line(slide, theme, fx, fy, tx, ty, "muted", 1.75)
        card_l = blk.lines[0] if len(blk.lines) > 0 else ""
        card_r = blk.lines[1] if len(blk.lines) > 1 else ""
        if card_l:
            lx = fx + (tx - fx) * 0.22 - Inches(0.2)
            ly = fy + (ty - fy) * 0.22 - Inches(0.15)
            add_text(slide, lx, ly, Inches(0.4), Inches(0.3), theme, card_l,
                     size=10, color_name="muted", bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if card_r:
            rx = fx + (tx - fx) * 0.78 - Inches(0.2)
            ry = fy + (ty - fy) * 0.78 - Inches(0.15)
            add_text(slide, rx, ry, Inches(0.4), Inches(0.3), theme, card_r,
                     size=10, color_name="muted", bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    for blk in entity_blocks:
        x, y, w, h = rects[blk.title]
        add_rect(slide, x, y, w, h, theme, "base_2", rounded=True)
        head_h = Inches(0.4)
        add_rect(slide, x, y, w, head_h, theme, "main", rounded=True)
        add_text(slide, x, y, w, head_h, theme, blk.title,
                 size=13, color_name="on_main", bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if blk.highlight:
            # カードは面積が大きくaccent塗りだとP2を超過しやすいため枠線のみで強調
            # （layered_stackと同じ判断）。
            _add_outline_rect(slide, theme, x, y, w, h, "accent")
        items = block_items(blk)
        if items:
            add_items_text(slide, x + Inches(0.14), y + head_h + Inches(0.08),
                            w - Inches(0.28), h - head_h - Inches(0.16), theme,
                            items, size=11, anchor=MSO_ANCHOR.TOP, bullet=(len(items) > 1))


R.register("layered_stack", render_layered_stack)
R.register("c4_context", render_c4_context)
R.register("sequence_diagram", render_sequence_diagram)
R.register("state_transition", render_state_transition)
R.register("er_diagram", render_er_diagram)
