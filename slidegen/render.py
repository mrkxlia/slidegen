"""
render.py — 内部データ構造 → ネイティブな PowerPoint 要素。

★最重要：設計ドキュメント §2-bis「PowerPointで編集可能であること」を厳守する。
  1. テキストは必ず text frame に入れる（画像化しない）
  2. 図形はネイティブシェイプ（Rectangle / Line）で作る
  3. 色・フォントは Theme 経由（potx提供後はテーマ参照に差し替え）
  4. 表は本物の PowerPoint テーブル（add_table）で作る
  5. すべて後から選択・移動・編集できる状態で置く

レイアウトは EMU ではなく Inches/Pt で素直に指定（16:9 = 13.333 x 7.5 inch）。
型ごとに render_<type>() を用意。新しい型はここに足すだけ（社内Claude Code向けの拡張点）。
"""
from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from .theme import Theme, DEFAULT_THEME, theme_from_potx
from .parser import Slide, split_emphasis

# 16:9 ワイド
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.7)          # §3 余白を確保（>0.5"）
CONTENT_W = SLIDE_W - MARGIN * 2


# ---------------------------------------------------------------------------
# 低レベルヘルパー（すべて編集可能なネイティブ要素を生成）
# ---------------------------------------------------------------------------
def _no_line(shape):
    shape.line.fill.background()

def _fill(shape, theme, color_name):
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme.rgb(color_name)
    _no_line(shape)

def add_rect(slide, x, y, w, h, theme, color_name, rounded=False):
    """矩形を描く。標準プリセット図形のみ使用（全環境で同一描画を保証）。
    rounded=False で直角、rounded=True で四隅角丸。
    部分角丸は使わない：角丸衝突は呼び出し側のレイアウト(重ね方)で回避する。"""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if rounded:
        # 角丸半径を控えめに固定（既定は大きすぎる）
        try:
            shp.adjustments[0] = 0.08
        except Exception:
            pass
    _fill(shp, theme, color_name)
    shp.shadow.inherit = False           # 影なし（§3 装飾を減らす）
    return shp

def add_text(slide, x, y, w, h, theme, runs, *, size, color_name,
             bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=None):
    """runs: [(text, emphasized_bool)] または str。emphasized は accent 色＋太字で表現。"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(runs, str):
        runs = [(runs, False)]
    for text, em in runs:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.name = font or theme.font
        r.font.bold = bold or em
        r.font.color.rgb = theme.rgb("accent") if em else theme.rgb(color_name)
    return box

def add_hline(slide, x, y, w, theme, color_name="main", weight=2.0):
    """ネイティブの直線。タイトル下の主張区切り。EMU整数を保証。"""
    ln = slide.shapes.add_connector(2, int(x), int(y), int(x + w), int(y))
    ln.line.color.rgb = theme.rgb(color_name)
    ln.line.width = Pt(weight)
    return ln


# ---------------------------------------------------------------------------
# 共通ヘッダー（kicker / headline / 区切り線）— 全型で使う
# ---------------------------------------------------------------------------
def render_header(slide, data: Slide, theme: Theme) -> Emu:
    y = MARGIN
    if data.props.get("kicker"):
        add_text(slide, MARGIN, y, CONTENT_W, Inches(0.4), theme,
                 data.props["kicker"], size=theme.sz_kicker, color_name="muted", bold=True)
        y += Inches(0.45)
    if data.props.get("headline"):
        add_text(slide, MARGIN, y, CONTENT_W, Inches(0.9), theme,
                 split_emphasis(data.props["headline"]),
                 size=theme.sz_headline, color_name="ink", bold=True)
        y += Inches(0.95)
    add_hline(slide, MARGIN, y, CONTENT_W, theme, "main", 2.0)
    return y + Inches(0.25)

def render_foot(slide, data: Slide, theme: Theme):
    foot = data.props.get("foot") or data.props.get("source")
    if foot:
        add_text(slide, MARGIN, SLIDE_H - Inches(0.55), CONTENT_W, Inches(0.35), theme,
                 foot, size=theme.sz_foot, color_name="muted")


# ---------------------------------------------------------------------------
# 型1：compare（横並び比較）— 要素数(col数)からカラム幅を自動決定（§4 要素数パターン化）
# ---------------------------------------------------------------------------
def render_compare(slide, data: Slide, theme: Theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)

    n = len(data.blocks)
    if n == 0:
        return
    gap = Inches(0.3)
    col_w = (CONTENT_W - gap * (n - 1)) / n
    bottom = SLIDE_H - Inches(0.7)
    col_h = bottom - top

    for i, blk in enumerate(data.blocks):
        x = MARGIN + i * (col_w + gap)
        # カード地（ベースの濃淡）
        add_rect(slide, x, top, col_w, col_h, theme, "base_2", rounded=True)
        # カラム見出し（強調列のみ accent、他は main）
        head_h = Inches(0.55)
        head_color = "accent" if blk.highlight else "main"
        hb = add_rect(slide, x, top, col_w, head_h, theme, head_color, rounded=True)
        add_text(slide, x, top, col_w, head_h, theme, blk.title,
                 size=theme.sz_col_title, color_name="on_main", bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 行（ラベル左／値右、罫線はグレー）
        row_y = top + head_h + Inches(0.2)
        rows = blk.rows if blk.rows else [("", l) for l in blk.lines]
        if rows:
            row_h = (col_h - head_h - Inches(0.4)) / len(rows)
            for j, (label, value) in enumerate(rows):
                ry = row_y + j * row_h
                pad = Inches(0.25)
                if label:
                    add_text(slide, x + pad, ry, col_w/2 - pad, row_h, theme, label,
                             size=theme.sz_body, color_name="muted", anchor=MSO_ANCHOR.MIDDLE)
                val_color = "accent" if blk.highlight else "ink"
                add_text(slide, x + col_w/2, ry, col_w/2 - pad, row_h, theme, value,
                         size=theme.sz_body, color_name=val_color, bold=True,
                         align=PP_ALIGN.RIGHT if label else PP_ALIGN.LEFT,
                         anchor=MSO_ANCHOR.MIDDLE)
                # 行間の罫線（最終行以外）
                if j < len(rows) - 1:
                    add_hline(slide, x + pad, ry + row_h, col_w - pad*2, theme, "rule", 1.0)


# ---------------------------------------------------------------------------
# 型2：kpi（数値ハイライト）— 大数字 1〜4 個
# ---------------------------------------------------------------------------
def render_kpi(slide, data: Slide, theme: Theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    n = len(data.blocks)
    if n == 0:
        return
    gap = Inches(0.3)
    col_w = (CONTENT_W - gap * (n - 1)) / n
    bottom = SLIDE_H - Inches(0.7)
    col_h = bottom - top
    for i, blk in enumerate(data.blocks):
        x = MARGIN + i * (col_w + gap)
        add_rect(slide, x, top, col_w, col_h, theme, "base_2", rounded=True)
        # 大数字（col title を数字として使う）
        add_text(slide, x, top + col_h*0.18, col_w, col_h*0.4, theme,
                 split_emphasis(blk.title),
                 size=theme.sz_stat, color_name=("accent" if blk.highlight else "main"),
                 bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # ラベル（lines / rows の最初）
        label = (blk.lines[0] if blk.lines else (blk.rows[0][1] if blk.rows else ""))
        add_text(slide, x, top + col_h*0.6, col_w, col_h*0.3, theme, label,
                 size=theme.sz_body, color_name="muted",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)


# ---------------------------------------------------------------------------
# 型3：process（プロセス／タイムライン）— ステップを横に並べ、矢印は控えめ（§3 矢印は補助）
# ---------------------------------------------------------------------------
def render_process(slide, data: Slide, theme: Theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    n = len(data.blocks)
    if n == 0:
        return
    gap = Inches(0.25)
    col_w = (CONTENT_W - gap * (n - 1)) / n
    bottom = SLIDE_H - Inches(0.9)
    box_h = bottom - top
    cy = top + box_h/2
    for i, blk in enumerate(data.blocks):
        x = MARGIN + i * (col_w + gap)
        color = "accent" if blk.highlight else "main"
        add_rect(slide, x, top, col_w, box_h, theme, "base_2", rounded=True)
        # ステップ番号バッジ
        badge = Inches(0.5)
        add_rect(slide, x + col_w/2 - badge/2, top + Inches(0.25), badge, badge, theme, color, rounded=True)
        add_text(slide, x + col_w/2 - badge/2, top + Inches(0.25), badge, badge, theme, str(i+1),
                 size=18, color_name="on_main", bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # タイトル
        add_text(slide, x + Inches(0.1), top + Inches(0.9), col_w - Inches(0.2), Inches(0.6), theme,
                 blk.title, size=theme.sz_col_title, color_name="ink", bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 説明（lines）— カード下半分の中央に配置
        if blk.lines:
            add_text(slide, x + Inches(0.15), top + Inches(1.5), col_w - Inches(0.3), box_h - Inches(1.7),
                     theme, " ".join(blk.lines), size=13, color_name="muted",
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # ステップ間の小さな矢印（補助・グレー）
        if i < n - 1:
            ax = x + col_w + gap/2
            tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, ax - Inches(0.07), cy - Inches(0.08),
                                         Inches(0.14), Inches(0.16))
            tri.rotation = 90
            _fill(tri, theme, "muted")
            tri.shadow.inherit = False


# ---------------------------------------------------------------------------
# ディスパッチ
# ---------------------------------------------------------------------------
# 型名→render関数の登録表。各モジュール(render_*.py)は import 時に
# register / register_many でここへ自分の型を足す（__init__.py が全 render_* を読む）。
RENDERERS: dict = {}

def register(name, fn):
    """型を1つ登録する。"""
    RENDERERS[name] = fn

def register_many(names, fn):
    """同一 render 関数を複数の型名（base + variants 等）へまとめて登録する。"""
    for name in names:
        RENDERERS[name] = fn

register("compare", render_compare)
register("kpi", render_kpi)
register("process", render_process)

def _strip_placeholders(slide):
    """白紙レイアウトに残るプレースホルダ(日付/ページ番号/タイトル枠など)を除去。
    これをしないと全スライドにマスター由来の枠が透けて表示される。"""
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)


def _strip_layout_placeholders(prs):
    """レイアウトとマスター側のプレースホルダも除去し、継承表示を止める。"""
    seen = set()
    for layout in prs.slide_layouts:
        for ph in list(layout.placeholders):
            ph._element.getparent().remove(ph._element)
        master = layout.slide_master
        if id(master) not in seen:
            seen.add(id(master))
            for ph in list(master.placeholders):
                ph._element.getparent().remove(ph._element)


def render_slide(prs, data: Slide, theme: Theme):
    layout = prs.slide_layouts[6]  # 6 = 白紙レイアウト（potx提供後は適切なレイアウトに）
    slide = prs.slides.add_slide(layout)
    _strip_placeholders(slide)     # 継承プレースホルダを除去（全ページ透け対策）
    fn = RENDERERS.get(data.type)
    if fn is None:
        raise ValueError(f"未対応の型: {data.type}（対応: {', '.join(RENDERERS)}）")
    fn(slide, data, theme)
    return slide


def build(slides_data, theme: Theme | None = None, template: str | None = None) -> Presentation:
    """
    template: 会社の .potx/.pptx パス。指定があればそれを土台にする（§2-bis ルール2）。
              未指定なら 16:9 の白紙プレゼンを作る。
    theme:    None（未指定）かつ template があれば potx のテーマ色を自動抽出する
              （§2-bis ルール3）。明示的に Theme を渡した場合はそれを優先する。
    """
    if template:
        prs = Presentation(template)
        if theme is None:
            theme = theme_from_potx(prs)   # template 提供かつ theme 未指定 → potx から自動抽出
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H
        _strip_layout_placeholders(prs)   # 白紙テンプレ由来の枠を全除去
    if theme is None:
        theme = DEFAULT_THEME
    for data in slides_data:
        render_slide(prs, data, theme)
    return prs
