"""
render_charts_shapes.py — 図形描画による「チャート風」の型群。

python-pptx のネイティブ Chart API では表現しにくい、または簡易に済ませたい型を
矩形/コネクタの標準プリセット図形の積み木で描く（waterfall・narrative_curve と同じ方針）。

- bullet          : 目標vs実績のバレットグラフ（ゲージの代替。箇条書きの bullets とは別の型）
- funnel          : 定量ファネル（段ごとに中央寄せで幅が減るバー。定性の段のみなら funnel_steps）
- football_field  : 評価レンジの横バー比較（M&Aのvaluation football field chart）
- harvey_ball_table : 定性比較表（●◐○）。OVALのリング＋PIEの部分塗りで4段階を表現
- marimekko         : 列幅=規模・縦100%積み上げのマリメッコチャート
- treemap           : 面積=構成比の決定的スライス&ダイス（1階層）
- sankey            : 左右2段の簡易フロー図。ノード=矩形、フロー=太さ可変の直線コネクタ

設計思想：標準プリセット図形のみ・回転/flip禁止。色は theme 経由。数値パースは
render_util.parse_number（typoは 0.0 として警告ログに残す）。marimekko/treemap の
highlight はセル面積が大きく accent 塗りだと P2（accent面積8%上限）に抵触しうるため、
アウトライン枠線＋周辺の識別ラベル（列名・合計等）の文字色で強調する（塗りは変えない。
セル内文字はコントラスト優先で背景色に応じた色のまま変えない）。
"""
from __future__ import annotations
import logging

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE

from . import render as R
from .render import (add_rect, add_text, add_hline, render_header, render_foot,
                     SLIDE_W, SLIDE_H, MARGIN, CONTENT_W)
from .parser import Slide
from .render_util import parse_number, columns_geometry

_log = logging.getLogger(__name__)


def _fmt_num(v: float) -> str:
    if abs(v - round(v)) < 1e-9:
        return f"{int(round(v)):,}"
    return f"{v:,.1f}"


def _row_num(b, i: int, default: float, context: str) -> float:
    """b.rows[i] があればその値、無ければ b.lines[i]、どちらも無ければ default。
    ラベル付き rows（`実績 "82"`）／無ラベル lines（`"82"`）のどちらの書き方でも
    位置引数として同じように拾う（ラベル文字列は表示用の自由文字列で意味は持たせない）。"""
    if i < len(b.rows):
        return parse_number(b.rows[i][1], context=context)
    if i < len(b.lines):
        return parse_number(b.lines[i], context=context)
    return default


def _unit_note(slide, theme, x, y, w, unit):
    if unit:
        add_text(slide, x, y, w, Inches(0.3), theme,
                 f"（単位：{unit}）", size=10, color_name="muted")


def _clamp_x(x, w):
    """テキストボックスがスライド外にはみ出さないよう左右をクランプ（narrative_curve と同じ発想）。"""
    return max(MARGIN * 0.3, min(x, SLIDE_W - MARGIN * 0.3 - w))


def _add_outline_rect(slide, theme, x, y, w, h, color_name, weight=2.0):
    """塗りなし・枠線のみの矩形（marimekko/treemap の highlight 用。線なので P2 accent算入対象外）。"""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), int(w), int(h))
    shp.fill.background()
    shp.line.color.rgb = theme.rgb(color_name)
    shp.line.width = Pt(weight)
    shp.shadow.inherit = False
    return shp


# ---------------------------------------------------------------------------
# bullet（目標vs実績のバレットグラフ）
# ---------------------------------------------------------------------------
def render_bullet(slide, data: Slide, theme):
    """col の行 = KPI 1個（上限4）。rows/lines の1行目=実績、2行目=目標、
    3行目(任意)=上限（省略時は実績・目標の1.15倍を自動設定）。
    記法:
      col "売上達成率" highlight
        実績 "82"
        目標 "100"
        上限 "120"
    """
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    rows = data.blocks[:4]
    n = len(rows)
    if n == 0:
        return

    unit = data.props.get("unit", "")
    _unit_note(slide, theme, MARGIN, top - Inches(0.05), CONTENT_W, unit)

    bottom = SLIDE_H - Inches(0.9)
    plot_top = top + Inches(0.25)
    row_h = (bottom - plot_top) / n

    title_w = Inches(2.3)
    value_w = Inches(1.5)
    track_x = MARGIN + title_w + Inches(0.15)
    track_w = CONTENT_W - title_w - value_w - Inches(0.3)
    track_h = Inches(0.32)

    for i, b in enumerate(rows):
        row_y = plot_top + i * row_h
        mid_y = row_y + row_h / 2

        actual = _row_num(b, 0, 0.0, "bullet")
        target = _row_num(b, 1, actual, "bullet")
        limit = _row_num(b, 2, 0.0, "bullet")
        if limit <= 0:
            limit = max(actual, target, 1.0) * 1.15

        add_text(slide, MARGIN, row_y, title_w - Inches(0.1), row_h, theme, b.title,
                 size=14, color_name="ink", bold=True, anchor=MSO_ANCHOR.MIDDLE)

        track_y = mid_y - track_h / 2
        add_rect(slide, track_x, track_y, track_w, track_h, theme, "base_2", rounded=False)

        frac_actual = max(0.0, min(actual / limit, 1.0))
        if frac_actual > 0:
            bar_h = track_h * 0.55
            bar_w = track_w * frac_actual
            add_rect(slide, track_x, mid_y - bar_h / 2, bar_w, bar_h, theme,
                     "accent" if b.highlight else "main", rounded=False)

        frac_target = max(0.0, min(target / limit, 1.0))
        marker_w = Inches(0.045)
        marker_h = track_h * 1.3
        marker_x = track_x + track_w * frac_target - marker_w / 2
        add_rect(slide, marker_x, mid_y - marker_h / 2, marker_w, marker_h, theme,
                 "ink", rounded=False)

        value_text = f"{_fmt_num(actual)} / {_fmt_num(target)}"
        add_text(slide, track_x + track_w + Inches(0.15), row_y, value_w, row_h, theme,
                 value_text, size=13, color_name="ink", bold=True,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------------------
# funnel（定量ファネル）
# ---------------------------------------------------------------------------
def render_funnel(slide, data: Slide, theme):
    """col の行 = 段1個（上限6）。1行目 = 値（数値）。
    記法:
      col "訪問"
        "10000"
      col "登録" highlight
        "1200"
    """
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    bars = data.blocks[:6]
    n = len(bars)
    if n == 0:
        return

    unit = data.props.get("unit", "")
    vals = [_row_num(b, 0, 0.0, "funnel") for b in bars]
    vmax = max(vals) if any(v > 0 for v in vals) else 1.0

    bottom = SLIDE_H - Inches(0.9)
    plot_top = top + Inches(0.15)
    plot_h = bottom - plot_top
    gap = Inches(0.16)
    max_bar_h = Inches(0.75)
    bar_h = min(max_bar_h, (plot_h - gap * (n - 1)) / n)
    total_h = bar_h * n + gap * (n - 1)
    start_y = plot_top + max(0, (plot_h - total_h) / 2)

    max_w = min(Inches(8.0), CONTENT_W)
    min_w = Inches(1.2)

    for i, b in enumerate(bars):
        frac = vals[i] / vmax if vmax else 0.0
        w = max(min_w, max_w * max(0.0, min(frac, 1.0)))
        x = MARGIN + (CONTENT_W - w) / 2
        y = start_y + i * (bar_h + gap)
        color = "accent" if b.highlight else "main"
        add_rect(slide, x, y, w, bar_h, theme, color, rounded=False)

        label = f"{b.title}　{_fmt_num(vals[i])}{unit}"
        if w >= Inches(2.0):
            add_text(slide, x, y, w, bar_h, theme, label, size=13,
                     color_name="on_main", bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        else:
            lx = x + w + Inches(0.1)
            lw = max(Inches(0.5), min(Inches(3.0), SLIDE_W - MARGIN - lx))
            add_text(slide, lx, y, lw, bar_h, theme, label, size=13,
                     color_name="ink", bold=True,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

        if i < n - 1 and vals[i]:
            rate = vals[i + 1] / vals[i] * 100
            ry = y + bar_h + gap / 2 - Inches(0.11)
            add_text(slide, SLIDE_W - MARGIN - Inches(1.3), ry, Inches(1.3), Inches(0.22),
                     theme, f"転換率 {rate:.0f}%", size=11, color_name="muted",
                     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------------------
# football_field（評価レンジの横バー比較）
# ---------------------------------------------------------------------------
def render_football_field(slide, data: Slide, theme):
    """col の行 = 評価手法1個（上限6）。1行 = "下限" "上限" の2値。
    任意 props `marker` で基準値（現在値等）の縦線を追加。
    記法:
      marker "100"
      col "DCF"
        "80" "120"
      col "類似会社比較" highlight
        "90" "140"
    """
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    rows = data.blocks[:6]
    n = len(rows)
    if n == 0:
        return

    unit = data.props.get("unit", "")
    _unit_note(slide, theme, MARGIN, top - Inches(0.05), CONTENT_W, unit)

    pairs = []
    for b in rows:
        vals = [parse_number(s, context="football_field") for s in b.lines[:2]]
        while len(vals) < 2:
            vals.append(0.0)
        lo, hi = vals[0], vals[1]
        if hi < lo:
            lo, hi = hi, lo
        pairs.append((lo, hi))

    marker_raw = data.props.get("marker")
    marker_val = parse_number(marker_raw, context="football_field") if marker_raw else None

    all_vals = [v for p in pairs for v in p] + ([marker_val] if marker_val is not None else [])
    vmin = min(all_vals) if all_vals else 0.0
    vmax = max(all_vals) if all_vals else 1.0
    if vmax <= vmin:
        vmax = vmin + 1.0
    pad = (vmax - vmin) * 0.1
    scale_min = vmin - pad
    scale_max = vmax + pad
    scale_range = scale_max - scale_min or 1.0

    label_w = Inches(2.0)
    track_x = MARGIN + label_w
    track_w = CONTENT_W - label_w

    def px(v):
        return track_x + track_w * (v - scale_min) / scale_range

    bottom = SLIDE_H - Inches(0.9)
    plot_top = top + Inches(0.3)
    plot_h = bottom - plot_top
    row_h = plot_h / n
    bar_h = min(Inches(0.4), row_h * 0.6)

    for i, (b, (lo, hi)) in enumerate(zip(rows, pairs)):
        y_center = plot_top + i * row_h + row_h / 2
        add_text(slide, MARGIN, y_center - Inches(0.2), label_w - Inches(0.15), Inches(0.4),
                 theme, b.title, size=13, color_name="ink", anchor=MSO_ANCHOR.MIDDLE)

        x1, x2 = px(lo), px(hi)
        bar_w = max(x2 - x1, Inches(0.03))
        color = "accent" if b.highlight else "main_2"
        add_rect(slide, x1, y_center - bar_h / 2, bar_w, bar_h, theme, color, rounded=False)

        lx1 = _clamp_x(x1 - Inches(0.95), Inches(0.9))
        add_text(slide, lx1, y_center - Inches(0.17), Inches(0.9), Inches(0.34), theme,
                 _fmt_num(lo), size=11, color_name="ink",
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        lx2 = _clamp_x(x2 + Inches(0.05), Inches(0.9))
        add_text(slide, lx2, y_center - Inches(0.17), Inches(0.9), Inches(0.34), theme,
                 _fmt_num(hi), size=11, color_name="ink",
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

        if i < n - 1:
            add_hline(slide, MARGIN, plot_top + (i + 1) * row_h, CONTENT_W, theme, "rule", 1.0)

    if marker_val is not None:
        mx = px(marker_val)
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(mx), int(plot_top),
                                          int(mx), int(bottom))
        conn.line.color.rgb = theme.rgb("accent")
        conn.line.width = Pt(1.5)
        conn.shadow.inherit = False
        mlx = _clamp_x(mx - Inches(0.4), Inches(0.8))
        add_text(slide, mlx, plot_top - Inches(0.3), Inches(0.8), Inches(0.25), theme,
                 _fmt_num(marker_val), size=10, color_name="accent", bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------------------
# harvey_ball_table（定性比較の●◐○表）
# ---------------------------------------------------------------------------
_HARVEY_LEVELS = (0, 25, 50, 75, 100)
_HARVEY_PIE_ADJ2 = {25: 0.0, 50: 54.0, 75: 108.0}   # adj1は常に162（12時起点・実機検証済み）


def _snap_harvey(v):
    v = max(0.0, min(100.0, v))
    return min(_HARVEY_LEVELS, key=lambda lvl: abs(lvl - v))


def _draw_harvey_ball(slide, theme, cx, cy, diameter, pct, fill_color):
    r = diameter / 2
    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(cx - r), int(cy - r), int(diameter), int(diameter))
    ring.shadow.inherit = False
    if pct >= 100:
        ring.fill.solid()
        ring.fill.fore_color.rgb = theme.rgb(fill_color)
        ring.line.fill.background()
    else:
        ring.fill.solid()
        ring.fill.fore_color.rgb = theme.rgb("base")
        ring.line.color.rgb = theme.rgb(fill_color)
        ring.line.width = Pt(1.25)
    if 0 < pct < 100:
        pie = slide.shapes.add_shape(MSO_SHAPE.PIE, int(cx - r), int(cy - r), int(diameter), int(diameter))
        pie.adjustments[0] = 162.0
        pie.adjustments[1] = _HARVEY_PIE_ADJ2[pct]
        pie.fill.solid()
        pie.fill.fore_color.rgb = theme.rgb(fill_color)
        pie.line.fill.background()
        pie.shadow.inherit = False


def render_harvey_ball_table(slide, data: Slide, theme):
    """comparison_matrix と同形（columns=評価軸、col=行）。値は 0/25/50/75/100 に丸めて
    Harveyボール（●◐○）で描く。行5×列5が上限（超過は切り捨て）。
    離散5値の定性記号であり、角度読み取り精度が問題になる円グラフの「❌」判断とは別物。
    記法:
      slide harvey_ball_table
        columns "機能" "価格" "サポート"
        col "自社" highlight
          "100" "75" "100"
    """
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)

    col_headers = data.props.get("columns_list")
    rows = data.blocks[:5]
    if not rows:
        return
    ncol = max((len(b.lines) for b in rows), default=0)
    if col_headers is None:
        single = data.props.get("columns")
        col_headers = [single] if single else [f"列{i+1}" for i in range(ncol)]
    col_headers = col_headers[:5]
    ncol = min(max(ncol, len(col_headers)), 5)
    if ncol == 0:
        return

    label_w = Inches(2.0)
    grid_w = CONTENT_W - label_w
    gap = Inches(0.06)
    cw = columns_geometry(grid_w, ncol, gap)

    bottom = SLIDE_H - Inches(0.7)
    avail_h = bottom - top
    header_h = Inches(0.5)
    nrow = len(rows)
    rh = (avail_h - header_h - gap * nrow) / nrow

    x0 = MARGIN + label_w
    for j in range(ncol):
        x = x0 + j * (cw + gap)
        add_rect(slide, x, top, cw, header_h, theme, "main", rounded=False)
        label = col_headers[j] if j < len(col_headers) else ""
        add_text(slide, x, top, cw, header_h, theme, label, size=12,
                 color_name="on_main", bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    diameter = Inches(0.35)
    for i, b in enumerate(rows):
        y = top + header_h + gap + i * (rh + gap)
        lab_color = "accent" if b.highlight else "base_2"
        add_rect(slide, MARGIN, y, label_w - gap, rh, theme, lab_color, rounded=False)
        add_text(slide, MARGIN + Inches(0.1), y, label_w - gap - Inches(0.2), rh, theme,
                 b.title, size=13, color_name="on_main" if b.highlight else "ink",
                 bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

        fill_color = "accent" if b.highlight else "main"
        for j in range(ncol):
            x = x0 + j * (cw + gap)
            cx = x + cw / 2
            cy = y + rh / 2
            raw = b.lines[j] if j < len(b.lines) else "0"
            pct = _snap_harvey(parse_number(raw, context="harvey_ball_table"))
            _draw_harvey_ball(slide, theme, cx, cy, diameter, pct, fill_color)


# ---------------------------------------------------------------------------
# marimekko（列幅=規模・縦100%積み上げ）
# ---------------------------------------------------------------------------
_SEGMENT_COLORS = ["main", "main_2", "main_3", "muted"]


def render_marimekko(slide, data: Slide, theme):
    """col = 列（上限5、幅は列合計に比例）。rows = (セグメント名, 値)（上限4、全列同順で
    書く）。セグメント色は行インデックスで [main, main_2, main_3, muted] に固定（全列で
    同じ色=同じセグメント）。highlight は列名・合計の文字色とアウトライン枠線のみで
    強調する（面積が大きく塗りを変えると accent 予算＝P2 を超過しやすいため）。
    記法:
      slide marimekko
        col "国内" highlight
          製品A "60"
          製品B "40"
        col "北米"
          製品A "30"
          製品B "50"
    """
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    cols = data.blocks[:5]
    n = len(cols)
    if n == 0:
        return

    unit = data.props.get("unit", "")
    _unit_note(slide, theme, MARGIN, top - Inches(0.05), CONTENT_W, unit)

    col_segments = []
    for b in cols:
        segs = []
        for i in range(4):
            if i < len(b.rows):
                segs.append(parse_number(b.rows[i][1], context="marimekko"))
            elif i < len(b.lines):
                segs.append(parse_number(b.lines[i], context="marimekko"))
        col_segments.append(segs)
    col_totals = [sum(s) for s in col_segments]
    grand_total = sum(col_totals) or 1.0

    bottom = SLIDE_H - Inches(0.9)
    plot_top = top + Inches(0.3)
    footer_h = Inches(0.6)
    plot_h = bottom - plot_top - footer_h
    gap = Inches(0.06)
    avail_w = CONTENT_W - gap * (n - 1)

    x = MARGIN
    for ci, b in enumerate(cols):
        total = col_totals[ci]
        w = avail_w * (total / grand_total) if grand_total > 0 else avail_w / n
        w = max(w, Inches(0.3))
        segs = col_segments[ci]
        y = plot_top
        for si, val in enumerate(segs):
            if total <= 0 or val <= 0:
                continue
            h = plot_h * (val / total)
            color = _SEGMENT_COLORS[si % len(_SEGMENT_COLORS)]
            add_rect(slide, x, y, w, h, theme, color, rounded=False)
            if h >= Inches(0.45) and w >= Inches(0.9):
                seg_label = b.rows[si][0] if si < len(b.rows) else ""
                txt_color = "on_main" if color in ("main", "main_2") else "ink"
                text = f"{seg_label} {_fmt_num(val)}" if seg_label else _fmt_num(val)
                add_text(slide, x + Inches(0.05), y, w - Inches(0.1), h, theme, text,
                         size=11, color_name=txt_color,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            y += h

        total_color = "accent" if b.highlight else "ink"
        add_text(slide, x, plot_top + plot_h + Inches(0.03), w, Inches(0.28), theme,
                 _fmt_num(total), size=12, color_name=total_color, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        add_text(slide, x, plot_top + plot_h + Inches(0.31), w, Inches(0.28), theme,
                 b.title, size=12, color_name=total_color, bold=b.highlight,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        if b.highlight:
            _add_outline_rect(slide, theme, x, plot_top, w, plot_h, "accent", weight=2.5)

        x += w + gap


# ---------------------------------------------------------------------------
# treemap（面積=構成比。決定的スライス&ダイス、1階層）
# ---------------------------------------------------------------------------
def render_treemap(slide, data: Slide, theme):
    """col = 項目（上限8）。1行目 = 値（数値）。DSL記述順のまま描画する（並べ替えない＝
    著者が大きい順に書けば大きい順に配置される決定的アルゴリズム）。横長の残余矩形は
    縦スライス、縦長なら横スライスを交互に適用する。highlight はアウトライン枠線のみで
    強調する（塗りは変えない。marimekko と同じ理由。セル文字はコントラスト優先で変えない）。
    記法:
      slide treemap
        col "クラウド" highlight
          "45"
        col "受託開発"
          "30"
    """
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    items = data.blocks[:8]
    n = len(items)
    if n == 0:
        return

    unit = data.props.get("unit", "")
    _unit_note(slide, theme, MARGIN, top - Inches(0.05), CONTENT_W, unit)

    vals = [_row_num(b, 0, 0.0, "treemap") for b in items]
    total = sum(vals) or 1.0

    bottom = SLIDE_H - Inches(0.7)
    plot_top = top + Inches(0.3)
    x, y, w, h = MARGIN, plot_top, CONTENT_W, bottom - plot_top
    gap = Inches(0.03)
    min_side = Inches(0.15)

    remaining_total = total
    for i, (b, val) in enumerate(zip(items, vals)):
        is_last = (i == n - 1)
        if is_last:
            cx, cy, cw, ch = x, y, w, h
        else:
            frac = (val / remaining_total) if remaining_total > 0 else (1.0 / (n - i))
            frac = max(0.0, min(frac, 1.0))
            if w >= h:
                cw = min(max(w * frac, min_side), w)
                cx, cy, ch = x, y, h
                x, w = x + cw + gap, w - cw - gap
            else:
                ch = min(max(h * frac, min_side), h)
                cx, cy, cw = x, y, w
                y, h = y + ch + gap, h - ch - gap
        remaining_total -= val

        color = _SEGMENT_COLORS[i % len(_SEGMENT_COLORS)]
        add_rect(slide, cx, cy, cw, ch, theme, color, rounded=False)
        if cw >= Inches(0.9) and ch >= Inches(0.5):
            txt_color = "on_main" if color in ("main", "main_2") else "ink"
            add_text(slide, cx + Inches(0.06), cy, cw - Inches(0.12), ch, theme,
                     f"{b.title} {_fmt_num(val)}", size=12, color_name=txt_color,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if b.highlight:
            _add_outline_rect(slide, theme, cx, cy, cw, ch, "accent", weight=2.5)


# ---------------------------------------------------------------------------
# sankey（左右2段の簡易フロー図）
# ---------------------------------------------------------------------------
_SANKEY_MAX_LEFT = 4
_SANKEY_MAX_RIGHT = 4
_SANKEY_MAX_FLOWS = 8


def _sankey_node_flows(b):
    """左ノード1個分のフロー先リスト [(右ノード名, 値), ...]。
    rows（`右ノード名 "値"`）があればそれを使う。rows が無い最小入力等のフォールバックでは
    lines を「先1」「先2」…というラベル無しフローとして扱う（クラッシュしない頑健性用）。"""
    if b.rows:
        return [(label, parse_number(val, context="sankey")) for label, val in b.rows]
    return [(f"先{i+1}", parse_number(v, context="sankey")) for i, v in enumerate(b.lines)]


def _sankey_node_heights(totals, avail_h, gap_v, min_h):
    n = len(totals)
    if n == 0:
        return []
    grand = sum(totals) or 1.0
    slot = avail_h - gap_v * (n - 1)
    heights = [max(slot * (t / grand) if grand > 0 else slot / n, min_h) for t in totals]
    total_h = sum(heights)
    if total_h > slot > 0:
        scale = slot / total_h
        heights = [h * scale for h in heights]
    return heights


def render_sankey(slide, data: Slide, theme):
    """col = 左ノード（上限4）。col 直下の各行「右ノード名 "値"」= そのノードからのフロー
    （右ノードは出現順に union、上限4。フロー総数は上限8で超過分は切り捨て）。
    ノードの高さ・フローの太さは値に比例する。左ノードの highlight はノード塗りと
    そこから出るフローの線色を accent にする。
    記法:
      slide sankey
        headline "流入チャネル別の転換フロー"
        unit "件"
        col "広告経由" highlight
          無料登録 "60"
          直接購入 "15"
        col "オーガニック"
          無料登録 "40"
    直線・テーパー無しのv1簡易版（カスタムジオメトリ禁止のため。narrative_curveと同じ判断）。
    """
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)

    left_nodes = data.blocks[:_SANKEY_MAX_LEFT]
    if not left_nodes:
        return

    raw_flows = []  # (left_idx, right_name, value)
    for li, b in enumerate(left_nodes):
        for rname, val in _sankey_node_flows(b):
            raw_flows.append((li, rname, val))

    right_order = []
    for _, rname, _ in raw_flows:
        if rname not in right_order and len(right_order) < _SANKEY_MAX_RIGHT:
            right_order.append(rname)

    flows = []
    dropped = 0
    for li, rname, val in raw_flows:
        if rname not in right_order or len(flows) >= _SANKEY_MAX_FLOWS:
            dropped += 1
            continue
        flows.append((li, right_order.index(rname), val, left_nodes[li].highlight))
    if dropped:
        _log.warning(
            "sankey: 上限（右ノード%d・フロー%d）超過のため%d件のフローを切り捨てました。",
            _SANKEY_MAX_RIGHT, _SANKEY_MAX_FLOWS, dropped)
    if not right_order or not flows:
        return

    unit = data.props.get("unit", "")
    _unit_note(slide, theme, MARGIN, top - Inches(0.05), CONTENT_W, unit)

    bottom = SLIDE_H - Inches(0.9)
    plot_top = top + Inches(0.3)
    plot_h = bottom - plot_top
    gap_v = Inches(0.22)
    min_h = Inches(0.3)
    label_w = Inches(1.8)
    node_w = Inches(0.35)

    left_x = MARGIN + label_w
    right_x = SLIDE_W - MARGIN - label_w - node_w

    left_totals = [0.0] * len(left_nodes)
    right_totals = [0.0] * len(right_order)
    for li, ri, val, _hl in flows:
        left_totals[li] += val
        right_totals[ri] += val

    left_heights = _sankey_node_heights(left_totals, plot_h, gap_v, min_h)
    right_heights = _sankey_node_heights(right_totals, plot_h, gap_v, min_h)

    left_y, y = [], plot_top
    for h in left_heights:
        left_y.append(y)
        y += h + gap_v
    right_y, y = [], plot_top
    for h in right_heights:
        right_y.append(y)
        y += h + gap_v

    # --- フロー（先に描き、ノードを後から重ねて端をきれいに見せる） ---
    max_val = max((f[2] for f in flows), default=1.0) or 1.0
    left_cursor = [0.0] * len(left_nodes)
    right_cursor = [0.0] * len(right_order)
    for li, ri, val, hl in flows:
        lt = left_totals[li] or 1.0
        rt = right_totals[ri] or 1.0
        ly_mid = left_y[li] + left_heights[li] * ((left_cursor[li] + val / 2) / lt)
        ry_mid = right_y[ri] + right_heights[ri] * ((right_cursor[ri] + val / 2) / rt)
        left_cursor[li] += val
        right_cursor[ri] += val

        weight = 2.0 + 22.0 * max(0.0, min(val / max_val, 1.0))
        color = "accent" if hl else "main_3"
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                          int(left_x + node_w), int(ly_mid),
                                          int(right_x), int(ry_mid))
        conn.line.color.rgb = theme.rgb(color)
        conn.line.width = Pt(weight)
        conn.shadow.inherit = False

    # --- ノード（左） ---
    for li, b in enumerate(left_nodes):
        color = "accent" if b.highlight else "main"
        add_rect(slide, left_x, left_y[li], node_w, left_heights[li], theme, color, rounded=False)
        label = f"{b.title}　{_fmt_num(left_totals[li])}{unit}"
        add_text(slide, MARGIN, left_y[li], label_w - Inches(0.1), left_heights[li], theme,
                 label, size=12, color_name="ink", align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # --- ノード（右） ---
    for ri, rname in enumerate(right_order):
        add_rect(slide, right_x, right_y[ri], node_w, right_heights[ri], theme, "main_2", rounded=False)
        label = f"{rname}　{_fmt_num(right_totals[ri])}{unit}"
        lx = right_x + node_w + Inches(0.1)
        add_text(slide, lx, right_y[ri], SLIDE_W - MARGIN - lx, right_heights[ri], theme,
                 label, size=12, color_name="ink", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


R.register("bullet", render_bullet)
R.register("funnel", render_funnel)
R.register("football_field", render_football_field)
R.register("harvey_ball_table", render_harvey_ball_table)
R.register("marimekko", render_marimekko)
R.register("treemap", render_treemap)
R.register("sankey", render_sankey)
