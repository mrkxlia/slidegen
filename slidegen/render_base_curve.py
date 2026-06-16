"""
render_base_curve.py — 基底レイアウト `narrative_curve` と variant ラッパー群。

折れ線（曲線）＋各点の注釈ピンで「推移と山場」を見せる基底。
感情曲線 / sparkline narrative / ストーリーの起伏 / 簡易トレンドを吸収する。

カスタムジオメトリ禁止のため、曲線は「点と点を直線セグメント(コネクタ)で繋ぐ」方式で描く。
各点にマーカー（小円）とラベルを置く。

DSL：各 col が1つの点。
  col "認知"          # title = 点のラベル（x軸）
    "+1"             # 1行目 = 高さ(-3..+3 程度の相対値)。无ければ0
  col "検討" highlight # highlight でその点を強調（山場）
    "-2"

variant:
  baseline : 中央に基準線を引くか
"""
from __future__ import annotations
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

from . import render as R
from .render import (add_rect, add_text, add_hline, render_header, render_foot,
                     SLIDE_W, SLIDE_H, MARGIN, CONTENT_W)
from .parser import Slide


VARIANTS = {
    "emotion_arc":   {"baseline": True},
    "story_curve":   {"baseline": True},
    "trend_line":    {"baseline": False},
    "sparkline_narrative": {"baseline": True},
}


def _resolve(data: Slide) -> dict:
    name = data.props.get("variant") or data.type
    return dict(VARIANTS.get(name, {"baseline": True}))


def _value(b):
    if b.lines:
        try:
            return float(b.lines[0].replace("+", ""))
        except Exception:
            return 0.0
    return 0.0


def _line_seg(slide, theme, x1, y1, x2, y2, color="main", weight=2.5):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      int(x1), int(y1), int(x2), int(y2))
    conn.line.color.rgb = theme.rgb(color)
    conn.line.width = Pt(weight)
    conn.shadow.inherit = False
    return conn


def render_narrative_curve(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    v = _resolve(data)
    pts = data.blocks
    n = len(pts)
    if n < 2:
        return

    bottom = SLIDE_H - Inches(1.0)
    plot_top = top + Inches(0.3)
    plot_h = bottom - plot_top
    mid_y = plot_top + plot_h / 2

    # x座標を均等配置（左右に余白）
    x0 = MARGIN + Inches(0.5)
    x1 = SLIDE_W - MARGIN - Inches(0.5)
    step = (x1 - x0) / (n - 1)

    vals = [_value(b) for b in pts]
    vmax = max(3.0, max(abs(x) for x in vals) or 1.0)

    def py(val):
        # +が上、-が下。中心からの相対
        return mid_y - (val / vmax) * (plot_h / 2 - Inches(0.4))

    coords = [(x0 + i * step, py(vals[i])) for i in range(n)]

    # 基準線
    if v.get("baseline"):
        add_hline(slide, x0, int(mid_y), int(x1 - x0), theme, "base_2", weight=1.0)

    # 折れ線セグメント
    for i in range(n - 1):
        _line_seg(slide, theme, coords[i][0], coords[i][1],
                  coords[i+1][0], coords[i+1][1], "main", 2.5)

    # 各点：マーカー＋ラベル
    for i, b in enumerate(pts):
        cxp, cyp = coords[i]
        accent = b.highlight
        r = Inches(0.13) if accent else Inches(0.09)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     int(cxp - r), int(cyp - r), int(r * 2), int(r * 2))
        dot.fill.solid(); dot.fill.fore_color.rgb = theme.rgb("accent" if accent else "main")
        dot.line.fill.background(); dot.shadow.inherit = False
        # ラベル（x軸名）は下に。境界からはみ出さないよう左右をクランプ
        lw = step
        lx = cxp - lw / 2
        lx = max(MARGIN * 0.3, min(lx, SLIDE_W - MARGIN * 0.3 - lw))
        add_text(slide, int(lx), int(bottom + Inches(0.05)),
                 int(lw), Inches(0.5), theme, b.title,
                 size=12, color_name="accent" if accent else "ink",
                 bold=accent, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        # 強調点には吹き出しメモ（2行目があれば）
        if accent and len(b.lines) > 1:
            note = b.lines[1]
            ny = cyp - Inches(0.75) if cyp > mid_y else cyp + Inches(0.25)
            nx = cxp - Inches(1.2)
            nx = max(MARGIN * 0.3, min(nx, SLIDE_W - MARGIN * 0.3 - Inches(2.4)))
            add_text(slide, int(nx), int(ny),
                     int(Inches(2.4)), Inches(0.5), theme, note,
                     size=12, color_name="accent", bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


R.RENDERERS["narrative_curve"] = render_narrative_curve
for _name in VARIANTS:
    R.RENDERERS[_name] = render_narrative_curve
