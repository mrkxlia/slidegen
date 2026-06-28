"""
render_relations.py — 39パターン(Cone社)から、要素間の関係を示す図解系の型を追加。

実装方針：
- 各型は data.blocks の数（col の数）でレイアウトを自動決定
- 強調は accent のみ／装飾は最小（影なし・罫線はグレー）
- すべてネイティブシェイプ（編集可能）

型一覧（このモジュール）:
  matrix    … 2x2マトリクス（ポジショニング・4象限）
  cycle     … サイクル（PDCA等）。3〜4要素は円配置、5+は四角配置
  pyramid   … ピラミッド型（上位ほど規模が小）
  tree      … ツリー図（中心1→子N）
  formula   … 数式型（A × B = C のように要素関係を式で示す）
  timeline  … 時系列年表（左→右）。26則「時系列は左から右」
  image_left… 画像（または図解領域）左 + テキスト右。26則「画像は必ず左」
"""
from __future__ import annotations
import math
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from . import render as R
from .render import (add_rect, add_text, add_hline, render_header, render_foot,
                     SLIDE_W, SLIDE_H, MARGIN, CONTENT_W)
from .parser import Slide, split_emphasis


def _fill(shape, theme, color_name):
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme.rgb(color_name)
    shape.line.fill.background()
    shape.shadow.inherit = False


def _add_oval(slide, x, y, w, h, theme, color_name):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    _fill(shp, theme, color_name)
    return shp


def _add_triangle(slide, x, y, w, h, theme, color_name):
    shp = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, x, y, w, h)
    _fill(shp, theme, color_name)
    return shp


# ---------------------------------------------------------------------------
# matrix — 2x2 マトリクス（4象限）
#   ・横軸/縦軸ラベルは props["x_axis"], props["y_axis"] か "横軸"/"縦軸"
#   ・blocks 1〜4個。block.title=象限名, block.lines=説明。highlightで強調象限。
#   ・象限の配置順は記法どおり（左上・右上・左下・右下）
# ---------------------------------------------------------------------------
def render_matrix(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    bottom = SLIDE_H - Inches(0.9)
    plot_top = top + Inches(0.1)
    plot_h = bottom - plot_top
    # 軸ラベル領域を確保
    axis_pad_l = Inches(0.5)   # 縦軸ラベル分
    axis_pad_b = Inches(0.4)   # 横軸ラベル分
    plot_x = MARGIN + axis_pad_l
    plot_w = CONTENT_W - axis_pad_l
    plot_h = plot_h - axis_pad_b
    cell_w = plot_w / 2
    cell_h = plot_h / 2

    # 4象限の枠（記法順に左上・右上・左下・右下）
    positions = [(0,0),(1,0),(0,1),(1,1)]
    for i, blk in enumerate(data.blocks[:4]):
        cx, cy = positions[i]
        x = plot_x + cx * cell_w
        y = plot_top + cy * cell_h
        bg = "main_3" if blk.highlight else "base_2"
        add_rect(slide, x + Inches(0.04), y + Inches(0.04),
                 cell_w - Inches(0.08), cell_h - Inches(0.08), theme, bg, rounded=True)
        title_color = "accent" if blk.highlight else "ink"
        add_text(slide, x + Inches(0.2), y + Inches(0.2), cell_w - Inches(0.4), Inches(0.5),
                 theme, split_emphasis(blk.title), size=16, color_name=title_color, bold=True)
        if blk.lines:
            add_text(slide, x + Inches(0.2), y + Inches(0.7), cell_w - Inches(0.4), cell_h - Inches(0.9),
                     theme, " ".join(blk.lines), size=12, color_name="muted")

    # 軸線（中心の十字）— ネイティブの直線
    cx0 = plot_x + cell_w
    cy0 = plot_top + cell_h
    add_hline(slide, plot_x, cy0, plot_w, theme, "main", 1.5)
    # 縦線も直接
    v = slide.shapes.add_connector(2, int(cx0), int(plot_top), int(cx0), int(plot_top + plot_h))
    v.line.color.rgb = theme.rgb("main"); v.line.width = Pt(1.5)

    # 軸ラベル
    x_axis = data.props.get("x_axis") or data.props.get("横軸") or ""
    y_axis = data.props.get("y_axis") or data.props.get("縦軸") or ""
    if x_axis:
        add_text(slide, plot_x, plot_top + plot_h + Inches(0.05), plot_w, axis_pad_b,
                 theme, x_axis, size=11, color_name="muted", align=PP_ALIGN.CENTER)
    if y_axis:
        # 簡略化：縦書きではなく左マージンに横書きで（編集性優先）
        add_text(slide, MARGIN, plot_top, axis_pad_l - Inches(0.1), plot_h, theme,
                 y_axis, size=11, color_name="muted", anchor=MSO_ANCHOR.MIDDLE,
                 align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# cycle — サイクル（PDCA等）
#   要素3〜4: 円形配置（円型サイクル）
#   要素5+  : 四角の輪を作る
# ---------------------------------------------------------------------------
def render_cycle(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    n = len(data.blocks)
    if n == 0:
        return
    bottom = SLIDE_H - Inches(0.7)
    cx = SLIDE_W / 2
    cy = (top + bottom) / 2
    radius = min((bottom - top), CONTENT_W) / 2 - Inches(0.8)
    node = Inches(1.4)

    for i, blk in enumerate(data.blocks):
        angle = -math.pi/2 + 2*math.pi * i / n  # 12時方向から時計回り
        ex = cx + radius * math.cos(angle) - node/2
        ey = cy + radius * math.sin(angle) - node/2
        color = "accent" if blk.highlight else "main"
        _add_oval(slide, ex, ey, node, node, theme, color)
        add_text(slide, ex, ey, node, node, theme, blk.title,
                 size=14, color_name="on_main", bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 説明は外側に置く（矢印を避けるためノードのさらに外側へ）
        if blk.lines:
            ox = cx + (radius + Inches(1.1)) * math.cos(angle) - Inches(1.4)
            oy = cy + (radius + Inches(1.1)) * math.sin(angle) - Inches(0.25)
            add_text(slide, ox, oy, Inches(2.8), Inches(0.5), theme, " ".join(blk.lines),
                     size=11, color_name="muted", align=PP_ALIGN.CENTER)
        # 矢印（次のノードへ。控えめなグレー）
        if n >= 2:
            angle_next = -math.pi/2 + 2*math.pi * ((i+1) % n) / n
            sx = cx + radius * math.cos(angle + 0.5) * 0.93
            sy = cy + radius * math.sin(angle + 0.5) * 0.93
            ex2 = cx + radius * math.cos(angle_next - 0.5) * 0.93
            ey2 = cy + radius * math.sin(angle_next - 0.5) * 0.93
            arr = slide.shapes.add_connector(2, int(sx), int(sy), int(ex2), int(ey2))
            arr.line.color.rgb = theme.rgb("muted")
            arr.line.width = Pt(1.5)


# ---------------------------------------------------------------------------
# pyramid — ピラミッド型（上位ほど規模が小）
#   blocks 3〜5 段。上から並べる。highlightは accent。
# ---------------------------------------------------------------------------
def render_pyramid(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    n = len(data.blocks)
    if n == 0:
        return
    bottom = SLIDE_H - Inches(0.7)
    pyr_h = bottom - top - Inches(0.2)
    layer_h = pyr_h / n
    max_w = Inches(7.0)
    cx = SLIDE_W / 2

    for i, blk in enumerate(data.blocks):
        # 台形を矩形で近似（編集容易さ優先）
        w = max_w * (i + 1) / n
        x = cx - w / 2
        y = top + Inches(0.1) + i * layer_h
        color = "accent" if blk.highlight else "main"
        add_rect(slide, x, y + Inches(0.05), w, layer_h - Inches(0.1), theme, color, rounded=False)
        add_text(slide, x, y, w, layer_h, theme, blk.title,
                 size=18, color_name="on_main", bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 右側に説明（スライド右端を超えないよう動的計算）
        if blk.lines:
            desc_x = x + w + Inches(0.3)
            desc_w = max(Inches(1.0), SLIDE_W - MARGIN - desc_x)
            add_text(slide, desc_x, y, desc_w, layer_h, theme,
                     " ".join(blk.lines), size=12, color_name="muted",
                     anchor=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------------------
# tree — ツリー図（中心ノード→子ノード横一列）
#   1個目を親、残りを子として扱う
# ---------------------------------------------------------------------------
def render_tree(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    if not data.blocks:
        return
    bottom = SLIDE_H - Inches(0.7)
    h_total = bottom - top
    parent = data.blocks[0]
    children = data.blocks[1:]

    node_w = Inches(2.2)
    node_h = Inches(0.9)
    px = SLIDE_W/2 - node_w/2
    py = top + Inches(0.2)
    color_p = "accent" if parent.highlight else "main"
    add_rect(slide, px, py, node_w, node_h, theme, color_p, rounded=True)
    add_text(slide, px, py, node_w, node_h, theme, parent.title,
             size=16, color_name="on_main", bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    if children:
        n = len(children)
        gap = Inches(0.3)
        cw = (CONTENT_W - gap*(n-1)) / n
        cy = py + node_h + Inches(0.9)
        for i, ch in enumerate(children):
            cx = MARGIN + i*(cw+gap)
            color = "accent" if ch.highlight else "main_2"
            add_rect(slide, cx, cy, cw, node_h, theme, color, rounded=True)
            add_text(slide, cx, cy, cw, node_h, theme, ch.title,
                     size=14, color_name="on_main", bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            # 親→子の接続線（グレー）
            line = slide.shapes.add_connector(2, int(px + node_w/2), int(py + node_h),
                                              int(cx + cw/2), int(cy))
            line.line.color.rgb = theme.rgb("rule"); line.line.width = Pt(1.5)
            # 子の説明
            if ch.lines:
                add_text(slide, cx, cy + node_h + Inches(0.15), cw, h_total - (cy + node_h + Inches(0.15) - top),
                         theme, " ".join(ch.lines), size=11, color_name="muted",
                         align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# formula — 数式型（A × B = C / A + B = C のように要素関係を式で示す）
#   props["operator"]: "x" (×) / "+" (+) / "=" の繰り返し（"x,x,=" など）
#   デフォルトは要素間 "x"、最後 "="。要素数 2〜4。
# ---------------------------------------------------------------------------
def render_formula(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    n = len(data.blocks)
    if n < 2:
        return
    # オペレータ列を組み立てる： n-1 個の演算子（最後の前は "="）
    op_spec = data.props.get("operator") or ("x," * (n-2) + "=") if n >= 2 else ""
    ops = [o.strip() for o in op_spec.split(",") if o.strip()] if "," in op_spec else list(op_spec)
    # デフォルト：最後の手前を "="、それ以前を "×"
    if not ops or len(ops) != n - 1:
        ops = ["x"] * (n - 2) + ["="] if n >= 2 else []
    op_disp = {"x": "×", "*": "×", "+": "＋", "=": "＝", "-": "−"}

    bottom = SLIDE_H - Inches(0.9)
    box_h = Inches(1.8)
    cy = (top + bottom) / 2
    # 要素ボックスと演算子テキストを並べる
    box_w = Inches(2.0)
    op_w  = Inches(0.6)
    total_w = box_w * n + op_w * (n - 1)
    start_x = SLIDE_W/2 - total_w/2

    x = start_x
    for i, blk in enumerate(data.blocks):
        color = "accent" if blk.highlight else "main"
        add_rect(slide, x, cy - box_h/2, box_w, box_h, theme, color, rounded=True)
        add_text(slide, x, cy - box_h/2, box_w, box_h*0.55, theme, blk.title,
                 size=18, color_name="on_main", bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
        if blk.lines:
            add_text(slide, x, cy, box_w, box_h*0.45, theme, blk.lines[0],
                     size=11, color_name="on_main",
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        x += box_w
        if i < n - 1:
            sym = op_disp.get(ops[i], ops[i])
            add_text(slide, x, cy - box_h/2, op_w, box_h, theme, sym,
                     size=36, color_name="ink", bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            x += op_w


# ---------------------------------------------------------------------------
# timeline — 時系列年表（左→右配置 / 26則「時系列は左から右」）
#   blocks 3〜6個。block.title=時点、block.lines=出来事
# ---------------------------------------------------------------------------
def render_timeline(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    n = len(data.blocks)
    if n == 0:
        return
    bottom = SLIDE_H - Inches(0.7)
    # 中央に横線を1本（main色）、上に時点ラベル、下にイベント説明（または交互）
    line_y = (top + bottom) / 2
    add_hline(slide, MARGIN, line_y, CONTENT_W, theme, "main", 2.5)
    step = CONTENT_W / (n - 1) if n > 1 else CONTENT_W
    for i, blk in enumerate(data.blocks):
        cx = MARGIN + (i * step if n > 1 else CONTENT_W/2)
        # ドット
        dot = Inches(0.35)
        color = "accent" if blk.highlight else "main"
        _add_oval(slide, cx - dot/2, line_y - dot/2, dot, dot, theme, color)
        # 上下交互配置（i偶数：上、奇数：下）
        up = (i % 2 == 0)
        label_h = Inches(0.5)
        desc_h = Inches(1.0)
        if up:
            ly = line_y - Inches(1.5)
            dy = line_y - Inches(0.55)
        else:
            ly = line_y + Inches(1.0)
            dy = line_y + Inches(0.25)
        # 時点（強）— ラベル幅を端で縮めて境界からはみ出さないようにする
        label_w = Inches(2.8)
        lx = cx - label_w/2
        # スライド境界にクランプ
        if lx < MARGIN:
            label_w = label_w - (MARGIN - lx)
            lx = MARGIN
        if lx + label_w > SLIDE_W - MARGIN:
            label_w = SLIDE_W - MARGIN - lx
        add_text(slide, lx, ly, label_w, label_h, theme,
                 split_emphasis(blk.title), size=14,
                 color_name=("accent" if blk.highlight else "ink"),
                 bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if blk.lines:
            add_text(slide, lx,
                     line_y - Inches(0.95) if up else line_y + Inches(0.55),
                     label_w, Inches(0.5), theme, " ".join(blk.lines),
                     size=11, color_name="muted",
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP if not up else MSO_ANCHOR.BOTTOM)


# ---------------------------------------------------------------------------
# image_left — 画像（または図解領域）左 + テキスト右（26則「画像は必ず左」）
#   実画像未指定時は左に main 色の図解プレースホルダーを置く
#   blocks の各要素を右側に縦に並べる（カード）
#   props["image"] にパスがあれば貼る
# ---------------------------------------------------------------------------
def render_image_left(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    bottom = SLIDE_H - Inches(0.7)
    # 左右分割
    gap = Inches(0.4)
    left_w = (CONTENT_W - gap) * 0.45
    right_w = CONTENT_W - gap - left_w
    lh = bottom - top
    img_path = data.props.get("image")
    if img_path:
        try:
            slide.shapes.add_picture(img_path, MARGIN, top, left_w, lh)
        except Exception:
            add_rect(slide, MARGIN, top, left_w, lh, theme, "main_3", rounded=True)
            add_text(slide, MARGIN, top, left_w, lh, theme,
                     "（画像を読み込めません）", size=14, color_name="muted",
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    else:
        # プレースホルダ（編集時に差し替え）
        add_rect(slide, MARGIN, top, left_w, lh, theme, "main_3", rounded=True)
        add_text(slide, MARGIN, top, left_w, lh, theme, "画像を配置",
                 size=18, color_name="main", bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 右側：blocks を縦カード
    items = data.blocks
    if items:
        rx = MARGIN + left_w + gap
        ih = (lh - Inches(0.2) * (len(items) - 1)) / len(items) if items else lh
        for i, blk in enumerate(items):
            y = top + i * (ih + Inches(0.2))
            color = "accent" if blk.highlight else "main"
            chip = Inches(0.35)
            add_rect(slide, rx, y + Inches(0.15), chip, chip, theme, color, rounded=True)
            add_text(slide, rx + chip + Inches(0.2), y, right_w - chip - Inches(0.2), Inches(0.5),
                     theme, split_emphasis(blk.title), size=16, color_name="ink", bold=True)
            if blk.lines:
                add_text(slide, rx + chip + Inches(0.2), y + Inches(0.5),
                         right_w - chip - Inches(0.2), ih - Inches(0.5), theme,
                         " ".join(blk.lines), size=12, color_name="muted")


# 登録
R.register("matrix", render_matrix)
R.register("cycle", render_cycle)
R.register("pyramid", render_pyramid)
R.register("tree", render_tree)
R.register("formula", render_formula)
R.register("timeline", render_timeline)
R.register("image_left", render_image_left)
