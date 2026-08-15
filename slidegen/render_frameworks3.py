"""
render_frameworks3.py — ビジネスフレーム個別型 第3弾（S5b）。

- vpc          : バリュープロポジションキャンバス（正式のsquare+circleは標準図形・回転なし
                 の制約下で再現できないため、左右2パネル×3段の矩形に簡略化）
- five_forces  : ファイブフォース（中央=業界内の競争、周辺4力、中心向き矢印）
- 3c           : 3C分析（顧客を頂点にした三角配置カード3枚。円の重なりは箇条書きが
                 入らないため venn2 では表現しない）
- bcg_matrix   : PPM（花形/問題児/金のなる木/負け犬の固定2x2。自由軸の matrix とは別型）
- empathy_map  : 共感マップ（上段2x2＝Think&Feel/See/Hear/Say&Do ＋ 下段 Pain/Gain の6ブロック）
- persona_card : ペルソナカード（左：写真プレースホルダ＋属性、右：ゴール等の縦積みカード）

設計思想：標準図形のみ。色は theme 経由。固定の意味論を持つので専用実装（前例：
render_frameworks.py の swot / render_frameworks2.py の bmc）。
"""
from __future__ import annotations
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from . import render as R
from .render import (add_rect, add_text, add_hline, render_header, render_foot,
                     SLIDE_W, SLIDE_H, MARGIN, CONTENT_W)
from .parser import Slide
from .render_util import block_items, add_items_text, fill_shape


def _cell(slide, theme, x, y, w, h, label, items, head_color, *, head_h, label_size, body_size):
    """帯（見出し）＋カード地＋箇条書き、の共通セル（bmc/swot の様式を汎用化）。"""
    add_rect(slide, int(x), int(y), int(w), int(h), theme, "base_2", rounded=True)
    add_rect(slide, int(x), int(y), int(w), int(head_h), theme, head_color, rounded=True)
    add_text(slide, int(x), int(y), int(w), int(head_h), theme, label,
              size=label_size, color_name="on_main", bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if items:
        add_items_text(slide, int(x + Inches(0.1)), int(y + head_h + Inches(0.06)),
                        int(w - Inches(0.2)), int(h - head_h - Inches(0.12)), theme,
                        items, size=body_size, anchor=MSO_ANCHOR.TOP, bullet=(len(items) > 1))


def _arrow(slide, theme, shape, x, y, w, h, color="muted"):
    shp = slide.shapes.add_shape(shape, int(x), int(y), int(w), int(h))
    fill_shape(shp, theme, color, no_shadow=True)
    return shp


def _connector(slide, theme, x1, y1, x2, y2, color="rule", weight=1.5):
    ln = slide.shapes.add_connector(2, int(x1), int(y1), int(x2), int(y2))
    ln.line.color.rgb = theme.rgb(color)
    ln.line.width = Pt(weight)
    return ln


# ---------------------------------------------------------------------------
# vpc — バリュープロポジションキャンバス
# 記法：タイトル無し col ×6 固定順
#   左上=Gain Creators / 左中=Products & Services / 左下=Pain Relievers
#   右上=Gains / 右中=Customer Jobs / 右下=Pains
# ---------------------------------------------------------------------------
_VPC_LEFT = ["Gain Creators｜ゲインクリエイター", "Products & Services｜製品・サービス",
             "Pain Relievers｜ペインリリーバー"]
_VPC_RIGHT = ["Gains｜顧客の便益", "Customer Jobs｜顧客の仕事", "Pains｜顧客の悩み"]


def render_vpc(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    blocks = data.blocks
    bottom = SLIDE_H - Inches(0.7)

    center_w = Inches(0.9)
    side_gap = Inches(0.2)
    panel_w = (CONTENT_W - center_w - 2 * side_gap) / 2
    left_x = MARGIN
    center_x = left_x + panel_w + side_gap
    right_x = center_x + center_w + side_gap

    panel_title_h = Inches(0.4)
    title_gap = Inches(0.08)
    sub_gap = Inches(0.08)
    sub_y0 = top + panel_title_h + title_gap
    sub_h = (bottom - sub_y0 - 2 * sub_gap) / 3

    add_rect(slide, int(left_x), int(top), int(panel_w), int(panel_title_h), theme, "main", rounded=True)
    add_text(slide, int(left_x), int(top), int(panel_w), int(panel_title_h), theme, "Value Map",
              size=12, color_name="on_main", bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, int(right_x), int(top), int(panel_w), int(panel_title_h), theme, "main_2", rounded=True)
    add_text(slide, int(right_x), int(top), int(panel_w), int(panel_title_h), theme, "Customer Profile",
              size=12, color_name="on_main", bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    for i in range(3):
        y = sub_y0 + i * (sub_h + sub_gap)
        litems = block_items(blocks[i]) if i < len(blocks) else []
        _cell(slide, theme, left_x, y, panel_w, sub_h, _VPC_LEFT[i], litems, "main",
              head_h=Inches(0.3), label_size=10, body_size=10)
        ritems = block_items(blocks[3 + i]) if 3 + i < len(blocks) else []
        _cell(slide, theme, right_x, y, panel_w, sub_h, _VPC_RIGHT[i], ritems, "main_2",
              head_h=Inches(0.3), label_size=10, body_size=10)

    # 中央：適合を示す双方向の矢印＋FIT
    cy = (top + bottom) / 2
    arrow_h = Inches(0.3)
    arrow_gap = Inches(0.3)
    start_y = cy - (arrow_h * 2 + arrow_gap) / 2
    _arrow(slide, theme, MSO_SHAPE.RIGHT_ARROW, center_x, start_y, center_w, arrow_h)
    _arrow(slide, theme, MSO_SHAPE.LEFT_ARROW, center_x, start_y + arrow_h + arrow_gap, center_w, arrow_h)
    add_text(slide, int(center_x), int(start_y + arrow_h), int(center_w), int(arrow_gap), theme,
              "FIT", size=9, color_name="muted", bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------------------
# five_forces — ファイブフォース分析
# 記法：タイトル無し col ×5 固定順（中央→上→左→右→下）
# ---------------------------------------------------------------------------
_FIVE_FORCES = [
    ("業界内の競争", "accent"),
    ("新規参入の脅威", "main"),
    ("売り手の交渉力", "main"),
    ("買い手の交渉力", "main"),
    ("代替品の脅威", "main"),
]


def render_five_forces(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    blocks = data.blocks
    bottom = SLIDE_H - Inches(0.7)
    avail_h = bottom - top

    center_w = Inches(3.0)
    arrow_w = Inches(0.5)
    gap_h = Inches(0.15)
    side_w = (CONTENT_W - center_w - 2 * arrow_w - 4 * gap_h) / 2
    gap_v = Inches(0.32)
    mid_h = avail_h * 0.42
    top_h = bottom_h = (avail_h - mid_h - 2 * gap_v) / 2

    x_left = MARGIN
    x_left_arrow = x_left + side_w + gap_h
    x_center = x_left_arrow + arrow_w + gap_h
    x_right_arrow = x_center + center_w + gap_h
    x_right = x_right_arrow + arrow_w + gap_h

    y_top = top
    y_mid = y_top + top_h + gap_v
    y_bottom = y_mid + mid_h + gap_v

    def items(i):
        return block_items(blocks[i]) if i < len(blocks) else []

    label, color = _FIVE_FORCES[0]
    _cell(slide, theme, x_center, y_mid, center_w, mid_h, label, items(0), color,
          head_h=Inches(0.36), label_size=12, body_size=10)
    label, color = _FIVE_FORCES[1]
    _cell(slide, theme, x_center, y_top, center_w, top_h, label, items(1), color,
          head_h=Inches(0.32), label_size=10, body_size=9)
    label, color = _FIVE_FORCES[2]
    _cell(slide, theme, x_left, y_mid, side_w, mid_h, label, items(2), color,
          head_h=Inches(0.32), label_size=10, body_size=9)
    label, color = _FIVE_FORCES[3]
    _cell(slide, theme, x_right, y_mid, side_w, mid_h, label, items(3), color,
          head_h=Inches(0.32), label_size=10, body_size=9)
    label, color = _FIVE_FORCES[4]
    _cell(slide, theme, x_center, y_bottom, center_w, bottom_h, label, items(4), color,
          head_h=Inches(0.32), label_size=10, body_size=9)

    av_w, av_h = Inches(0.4), Inches(0.22)
    _arrow(slide, theme, MSO_SHAPE.DOWN_ARROW,
           x_center + center_w / 2 - av_w / 2, y_top + top_h + (gap_v - av_h) / 2, av_w, av_h)
    _arrow(slide, theme, MSO_SHAPE.UP_ARROW,
           x_center + center_w / 2 - av_w / 2, y_mid + mid_h + (gap_v - av_h) / 2, av_w, av_h)
    ah_h = Inches(0.35)
    _arrow(slide, theme, MSO_SHAPE.RIGHT_ARROW, x_left_arrow, y_mid + mid_h / 2 - ah_h / 2, arrow_w, ah_h)
    _arrow(slide, theme, MSO_SHAPE.LEFT_ARROW, x_right_arrow, y_mid + mid_h / 2 - ah_h / 2, arrow_w, ah_h)


# ---------------------------------------------------------------------------
# 3c — 3C分析（顧客を頂点にした三角配置カード3枚）
# 記法：タイトル無し col ×3 固定順（顧客→自社→競合）
# ---------------------------------------------------------------------------
_3C = [("Customer｜市場・顧客", "main"), ("Company｜自社", "main_2"), ("Competitor｜競合", "muted")]


def render_three_c(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    blocks = data.blocks
    bottom = SLIDE_H - Inches(0.7)
    avail_h = bottom - top

    gap = Inches(0.3)
    top_card_w = Inches(5.2)
    row1_h = avail_h * 0.42
    row2_h = avail_h - row1_h - gap
    bottom_card_w = (CONTENT_W - gap) / 2

    top_x = MARGIN + (CONTENT_W - top_card_w) / 2
    top_y = top
    bl_x = MARGIN
    bl_y = top + row1_h + gap
    br_x = MARGIN + bottom_card_w + gap
    br_y = bl_y

    # 接続線を先に描き、カードの下に回す
    top_cx = top_x + top_card_w / 2
    top_bottom_y = top_y + row1_h
    bl_cx = bl_x + bottom_card_w / 2
    br_cx = br_x + bottom_card_w / 2
    _connector(slide, theme, top_cx, top_bottom_y, bl_cx, bl_y)
    _connector(slide, theme, top_cx, top_bottom_y, br_cx, br_y)
    _connector(slide, theme, bl_x + bottom_card_w, bl_y + row2_h / 2, br_x, br_y + row2_h / 2)

    def items(i):
        return block_items(blocks[i]) if i < len(blocks) else []

    label, color = _3C[0]
    _cell(slide, theme, top_x, top_y, top_card_w, row1_h, label, items(0), color,
          head_h=Inches(0.4), label_size=14, body_size=12)
    label, color = _3C[1]
    _cell(slide, theme, bl_x, bl_y, bottom_card_w, row2_h, label, items(1), color,
          head_h=Inches(0.4), label_size=14, body_size=12)
    label, color = _3C[2]
    _cell(slide, theme, br_x, br_y, bottom_card_w, row2_h, label, items(2), color,
          head_h=Inches(0.4), label_size=14, body_size=12)


# ---------------------------------------------------------------------------
# bcg_matrix — PPM（花形/問題児/金のなる木/負け犬の固定2x2。自由軸は matrix 型を使う）
# 記法：タイトル無し col ×4 固定順（花形→問題児→金のなる木→負け犬）
# ---------------------------------------------------------------------------
_BCG = [
    ("Star｜花形", "main"),
    ("Question Mark｜問題児", "accent"),
    ("Cash Cow｜金のなる木", "main_2"),
    ("Dog｜負け犬", "muted"),
]
_BCG_POSITIONS = [(0, 0), (1, 0), (0, 1), (1, 1)]


def render_bcg_matrix(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    blocks = data.blocks
    bottom = SLIDE_H - Inches(0.7)

    axis_pad_l = Inches(0.45)
    axis_pad_b = Inches(0.35)
    plot_top = top + Inches(0.05)
    plot_x = MARGIN + axis_pad_l
    plot_w = CONTENT_W - axis_pad_l
    plot_h = bottom - plot_top - axis_pad_b
    gap = Inches(0.2)
    cw = (plot_w - gap) / 2
    ch = (plot_h - gap) / 2

    for i, (c, r) in enumerate(_BCG_POSITIONS):
        x = plot_x + c * (cw + gap)
        y = plot_top + r * (ch + gap)
        label, color = _BCG[i]
        items = block_items(blocks[i]) if i < len(blocks) else []
        _cell(slide, theme, x, y, cw, ch, label, items, color,
              head_h=Inches(0.4), label_size=13, body_size=11)

    add_text(slide, int(plot_x), int(plot_top + plot_h + Inches(0.05)), int(plot_w), int(axis_pad_b),
              theme, "相対マーケットシェア（高 ← → 低）", size=10, color_name="muted",
              align=PP_ALIGN.CENTER)
    add_text(slide, int(MARGIN), int(plot_top), int(axis_pad_l - Inches(0.05)), int(plot_h),
              theme, "市場成長率（高↑）", size=10, color_name="muted",
              anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# empathy_map — 共感マップ（上段2x2＋下段 Pain/Gain の6ブロック）
# 記法：タイトル無し col ×6 固定順
#   考え・感情→見えているもの→聞いていること→発言・行動→Pain→Gain
# ---------------------------------------------------------------------------
_EMPATHY_UPPER = [
    ("Think & Feel｜考え・感情", "main"),
    ("See｜見えているもの", "main"),
    ("Hear｜聞いていること", "main"),
    ("Say & Do｜発言・行動", "main"),
]
_EMPATHY_UPPER_POS = [(0, 0), (1, 0), (0, 1), (1, 1)]
_EMPATHY_LOWER = [("Pain｜痛み・不満", "muted"), ("Gain｜望む成果", "main_2")]


def render_empathy_map(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    blocks = data.blocks
    bottom = SLIDE_H - Inches(0.7)
    avail_h = bottom - top

    gap = Inches(0.15)
    upper_h = avail_h * 0.62
    lower_h = avail_h - upper_h - gap
    cw = (CONTENT_W - gap) / 2
    uch = (upper_h - gap) / 2

    for i, (c, r) in enumerate(_EMPATHY_UPPER_POS):
        x = MARGIN + c * (cw + gap)
        y = top + r * (uch + gap)
        label, color = _EMPATHY_UPPER[i]
        items = block_items(blocks[i]) if i < len(blocks) else []
        _cell(slide, theme, x, y, cw, uch, label, items, color,
              head_h=Inches(0.34), label_size=11, body_size=10)

    lower_y = top + upper_h + gap
    for j in range(2):
        x = MARGIN + j * (cw + gap)
        label, color = _EMPATHY_LOWER[j]
        idx = 4 + j
        items = block_items(blocks[idx]) if idx < len(blocks) else []
        _cell(slide, theme, x, lower_y, cw, lower_h, label, items, color,
              head_h=Inches(0.34), label_size=12, body_size=10)


# ---------------------------------------------------------------------------
# persona_card — ペルソナカード
# 記法：
#   name "田中 花子"   role "経理マネージャー（38）"
#   col "プロフィール"        # rows: ラベル "値"（属性一覧）
#     経験 "経理歴12年"
#   col "ゴール"              # lines: 箇条書き
#     "月次決算を5営業日で締める"
#   col "課題" highlight
#     "紙のレシート照合に月40時間"
# ---------------------------------------------------------------------------
def render_persona_card(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    blocks = data.blocks
    bottom = SLIDE_H - Inches(0.7)
    avail_h = bottom - top

    gap = Inches(0.25)
    left_w = CONTENT_W * 0.32
    right_w = CONTENT_W - left_w - gap
    left_x = MARGIN
    right_x = left_x + left_w + gap

    add_rect(slide, int(left_x), int(top), int(left_w), int(avail_h), theme, "base_2", rounded=True)

    name = data.props.get("name", "")
    role = data.props.get("role", "")
    photo_d = Inches(1.1)
    photo_x = left_x + (left_w - photo_d) / 2
    photo_y = top + Inches(0.25)
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(photo_x), int(photo_y), int(photo_d), int(photo_d))
    fill_shape(oval, theme, "main_3", no_shadow=True)
    add_text(slide, int(photo_x), int(photo_y), int(photo_d), int(photo_d), theme,
              name[:1] if name else "?", size=28, color_name="main", bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    name_y = photo_y + photo_d + Inches(0.12)
    add_text(slide, int(left_x + Inches(0.1)), int(name_y), int(left_w - Inches(0.2)), Inches(0.4),
              theme, name, size=18, color_name="ink", bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    role_y = name_y + Inches(0.4)
    add_text(slide, int(left_x + Inches(0.1)), int(role_y), int(left_w - Inches(0.2)), Inches(0.55),
              theme, role, size=11, color_name="muted",
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    rule_y = role_y + Inches(0.55)
    add_hline(slide, int(left_x + Inches(0.15)), int(rule_y), int(left_w - Inches(0.3)), theme, "rule", 1.0)
    profile_y = rule_y + Inches(0.12)
    profile_h = bottom - profile_y - Inches(0.12)
    if blocks:
        items = block_items(blocks[0])
        if items:
            add_items_text(slide, int(left_x + Inches(0.15)), int(profile_y),
                            int(left_w - Inches(0.3)), int(profile_h), theme,
                            items, size=11, anchor=MSO_ANCHOR.TOP, bullet=(len(items) > 1))

    rest = blocks[1:]
    n = len(rest)
    if n:
        vgap = Inches(0.15)
        card_h = (avail_h - vgap * (n - 1)) / n
        for i, blk in enumerate(rest):
            y = top + i * (card_h + vgap)
            head_color = "accent" if blk.highlight else "main"
            _cell(slide, theme, right_x, y, right_w, card_h, blk.title, block_items(blk), head_color,
                  head_h=Inches(0.38), label_size=13, body_size=12)


# ---------------------------------------------------------------------------
# speaker_intro_card — 登壇者紹介カード（persona_card のOVAL写真意匠を単一フォーカスへ
#   簡略化。右側カード群は持たず、写真+名前+役職+bio箇条書きのみの中央寄せ構成）
# ---------------------------------------------------------------------------
def render_speaker_intro_card(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)
    bottom = SLIDE_H - Inches(0.7)

    name = data.props.get("name", "")
    role = data.props.get("role", "")
    cx = SLIDE_W / 2
    photo_d = Inches(1.8)
    photo_x = cx - photo_d / 2
    photo_y = top + Inches(0.2)
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(photo_x), int(photo_y), int(photo_d), int(photo_d))
    fill_shape(oval, theme, "main_3", no_shadow=True)
    add_text(slide, int(photo_x), int(photo_y), int(photo_d), int(photo_d), theme,
              name[:1] if name else "?", size=48, color_name="main", bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    name_y = photo_y + photo_d + Inches(0.2)
    add_text(slide, int(MARGIN), int(name_y), int(CONTENT_W), Inches(0.5), theme, name,
              size=24, color_name="ink", bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    role_y = name_y + Inches(0.5)
    add_text(slide, int(MARGIN), int(role_y), int(CONTENT_W), Inches(0.4), theme, role,
              size=14, color_name="muted",
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    bio_y = role_y + Inches(0.55)
    if data.blocks:
        items = block_items(data.blocks[0])
        if items:
            bio_w = CONTENT_W * 0.6
            bio_x = MARGIN + (CONTENT_W - bio_w) / 2
            add_items_text(slide, int(bio_x), int(bio_y), int(bio_w), int(bottom - bio_y),
                            theme, items, size=13, anchor=MSO_ANCHOR.TOP, bullet=(len(items) > 1))


R.register("vpc", render_vpc)
R.register("five_forces", render_five_forces)
R.register("3c", render_three_c)
R.register("bcg_matrix", render_bcg_matrix)
R.register("empathy_map", render_empathy_map)
R.register("persona_card", render_persona_card)
R.register("speaker_intro_card", render_speaker_intro_card)
