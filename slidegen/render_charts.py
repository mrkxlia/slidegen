"""
render_charts.py — ネイティブ編集可能チャート型。

python-pptx の Chart API で「本物のグラフ」を生成する（画像化しない＝編集可能）。
設計思想：円グラフ・3D・ゲージは作らない。棒/横棒/折れ線/積み上げ/100%積み上げ/クラスター/
面/散布図/バブルを提供。系列色は theme の main 系＋accent に限定（70:25:5 を守る）。

DSL（categories系。既存パーサで書ける形。render_chart が担当）:
  slide bar_chart
    headline "四半期売上"
    categories "Q1" "Q2" "Q3" "Q4"   # x軸ラベル（多値プロパティ→categories_list）
    unit "百万円"                      # y軸の単位（任意）
    col "売上"                         # col.title = 系列名
      "120"                           # 各 line = 数値（カテゴリ順）
      "150"
      "135"
      "180"
    col "目標"                         # 2系列目（clustered/lineで複数系列に）
      "100"
      "140"
      "160"
      "170"

型 → チャート種別（categories系）:
  bar_chart        : 縦棒（単一/複数系列=clustered）
  bar_horizontal   : 横棒
  line_chart       : 折れ線
  stacked_bar      : 積み上げ縦棒
  stacked_100_bar  : 100%積み上げ（円グラフ代替）
  clustered_bar    : クラスター縦棒（明示）
  area_chart       : 積み上げ面グラフ（推移の累積を面で示す）

DSL（x,y座標系。render_xy_chart が担当。categories の代わりに col の各行が1点）:
  slide scatter
    headline "広告費と売上の相関"
    x_label "広告費（百万円）"          # 任意：軸ラベル注記
    y_label "売上（億円）"
    col "既存店"                       # col = 系列（複数可）
      "10" "1.2"                      # 1行 = "x" "y" の2値
      "15" "1.8"
  slide bubble                        # 同形。1行 = "x" "y" "規模" の3値
    col "事業ポートフォリオ"
      "12" "30" "45"

型 → チャート種別（x,y座標系）:
  scatter : 散布図（2変数の相関）
  bubble  : バブルチャート（3変数目=バブルの大きさ）
"""
from __future__ import annotations
import logging

from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_MARKER_STYLE

from . import render as R
from .render import add_text, render_header, render_foot, SLIDE_H, MARGIN, CONTENT_W
from .parser import Slide

_log = logging.getLogger(__name__)


CHART_TYPES = {
    "bar_chart":       XL_CHART_TYPE.COLUMN_CLUSTERED,
    "clustered_bar":   XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar_horizontal":  XL_CHART_TYPE.BAR_CLUSTERED,
    "line_chart":      XL_CHART_TYPE.LINE_MARKERS,
    "stacked_bar":     XL_CHART_TYPE.COLUMN_STACKED,
    "stacked_100_bar": XL_CHART_TYPE.COLUMN_STACKED_100,
    "area_chart":      XL_CHART_TYPE.AREA_STACKED,
}

# x,y座標系チャート（categories を使わず col の各行を点として描く）
XY_CHART_TYPES = {
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "bubble":  XL_CHART_TYPE.BUBBLE,
}
_XY_STRIDE = {"scatter": 2, "bubble": 3}       # 1点あたりの値の個数（x,y[,規模]）
_XY_POINT_LIMIT = {"scatter": 20, "bubble": 10}  # 系列あたりの点数上限（超過は切り捨て）

# 系列に割り当てる色（main系＋accent。多すぎる系列は非推奨）
_SERIES_COLORS = ["main", "main_2", "main_3", "accent", "muted"]


def _nums(lines):
    out = []
    for s in lines:
        try:
            out.append(float(s.replace(",", "").replace("+", "")))
        except (ValueError, AttributeError):
            _log.warning("チャートの数値として解釈できない値 %r を 0.0 として扱います。", s)
            out.append(0.0)
    return out


def _style_legend(chart, theme, n_series):
    """凡例：系列が複数のときだけ表示（categories系・x,y座標系で共用）。"""
    if n_series > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(11)
        chart.legend.font.name = theme.font
    else:
        chart.has_legend = False


def _style_axes(chart, theme):
    """軸フォント・目盛（categories系・x,y座標系で共用。x,y座標系でも category_axis はx軸を返す）。"""
    try:
        cax = chart.category_axis
        cax.tick_labels.font.size = Pt(11)
        cax.tick_labels.font.name = theme.font
        vax = chart.value_axis
        vax.tick_labels.font.size = Pt(10)
        vax.tick_labels.font.name = theme.font
        vax.has_major_gridlines = True
    except (AttributeError, ValueError) as e:
        _log.warning("チャートの軸スタイル設定に失敗しました（%s）。既定の軸表示で続行します。", e)


def render_chart(slide, data: Slide, theme):
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)

    chart_type = CHART_TYPES.get(data.type, XL_CHART_TYPE.COLUMN_CLUSTERED)
    series = data.blocks
    if not series:
        return

    # カテゴリ（x軸）
    cats = data.props.get("categories_list")
    if cats is None:
        single = data.props.get("categories")
        # 無ければ系列の値数から連番
        ncat = max(len(b.lines) for b in series)
        cats = [single] if single else [f"{i+1}" for i in range(ncat)]

    cd = CategoryChartData()
    cd.categories = cats
    for b in series:
        vals = _nums(b.lines)
        # カテゴリ数に長さを合わせる
        while len(vals) < len(cats):
            vals.append(0.0)
        cd.add_series(b.title or "系列", vals[:len(cats)])

    # チャート配置領域（ヘッダー下〜フッタ上）
    x = MARGIN
    y = top + Inches(0.1)
    w = CONTENT_W
    h = SLIDE_H - Inches(0.9) - y

    gframe = slide.shapes.add_chart(chart_type, int(x), int(y), int(w), int(h), cd)
    chart = gframe.chart

    # --- スタイリング（装飾を抑える） ---
    chart.has_title = False
    _style_legend(chart, theme, len(series))

    # 系列色
    for i, plot_series in enumerate(chart.series):
        color = theme.rgb(_SERIES_COLORS[i % len(_SERIES_COLORS)])
        fmt = plot_series.format
        if data.type == "line_chart":
            fmt.line.color.rgb = color
            fmt.line.width = Pt(2.5)
        else:
            fmt.fill.solid()
            fmt.fill.fore_color.rgb = color

    # データラベル（単一系列の棒なら値を直接表示＝凡例不要の直接ラベリング）
    try:
        plot = chart.plots[0]
        if len(series) == 1 and data.type in ("bar_chart", "bar_horizontal", "clustered_bar"):
            plot.has_data_labels = True
            plot.data_labels.font.size = Pt(11)
            plot.data_labels.font.name = theme.font
            plot.data_labels.number_format = '#,##0'
            plot.data_labels.number_format_is_linked = False
    except (AttributeError, IndexError, KeyError) as e:
        _log.warning("チャートのデータラベル設定に失敗しました（%s）。ラベル無しで続行します。", e)

    # 軸フォント
    _style_axes(chart, theme)

    # 単位ラベル（右上に小さく）
    unit = data.props.get("unit")
    if unit:
        add_text(slide, x, top - Inches(0.05), w, Inches(0.3), theme,
                 f"（単位：{unit}）", size=10, color_name="muted")


def render_xy_chart(slide, data: Slide, theme):
    """散布図・バブルチャート。categories を使わず、col の各行を (x, y[, 規模]) の点として描く。"""
    top = render_header(slide, data, theme)
    render_foot(slide, data, theme)

    chart_type = XY_CHART_TYPES.get(data.type, XL_CHART_TYPE.XY_SCATTER)
    series = data.blocks
    if not series:
        return

    is_bubble = data.type == "bubble"
    stride = _XY_STRIDE.get(data.type, 2)
    limit = _XY_POINT_LIMIT.get(data.type, 20)

    cd = BubbleChartData() if is_bubble else XyChartData()
    for b in series:
        vals = _nums(b.lines)
        if len(vals) % stride:
            _log.warning(
                "チャート系列 %r の値の個数が %d の倍数ではありません。末尾の端数は切り捨てます。",
                b.title or "系列", stride,
            )
        n_points = min(len(vals) // stride, limit)
        s = cd.add_series(b.title or "系列")
        for i in range(n_points):
            px, py = vals[i * stride], vals[i * stride + 1]
            if is_bubble:
                s.add_data_point(px, py, vals[i * stride + 2])
            else:
                s.add_data_point(px, py)

    # チャート配置領域（ヘッダー下〜フッタ上）
    x = MARGIN
    y = top + Inches(0.1)
    w = CONTENT_W
    h = SLIDE_H - Inches(0.9) - y

    gframe = slide.shapes.add_chart(chart_type, int(x), int(y), int(w), int(h), cd)
    chart = gframe.chart

    chart.has_title = False
    _style_legend(chart, theme, len(series))

    # 系列色：scatter はマーカー、bubble はバブル本体の塗り
    for i, plot_series in enumerate(chart.series):
        color = theme.rgb(_SERIES_COLORS[i % len(_SERIES_COLORS)])
        try:
            if is_bubble:
                fmt = plot_series.format
                fmt.fill.solid()
                fmt.fill.fore_color.rgb = color
                fmt.line.fill.background()
            else:
                marker = plot_series.marker
                marker.style = XL_MARKER_STYLE.CIRCLE
                marker.size = 8
                marker.format.fill.solid()
                marker.format.fill.fore_color.rgb = color
                marker.format.line.fill.background()
        except (AttributeError, ValueError) as e:
            _log.warning("系列の装飾に失敗しました（%s）。既定の見た目で続行します。", e)

    if is_bubble:
        try:
            chart.plots[0].bubble_scale = 60
        except (AttributeError, IndexError, ValueError) as e:
            _log.warning("バブルサイズの調整に失敗しました（%s）。既定のスケールで続行します。", e)

    _style_axes(chart, theme)

    # 軸ラベル注記（右上に小さく。unit の代わりに x_label/y_label）
    x_label = data.props.get("x_label")
    y_label = data.props.get("y_label")
    parts = [p for p in (
        f"横軸：{x_label}" if x_label else None,
        f"縦軸：{y_label}" if y_label else None,
    ) if p]
    if parts:
        add_text(slide, x, top - Inches(0.05), w, Inches(0.3), theme,
                 f"（{'／'.join(parts)}）", size=10, color_name="muted")


R.register_many(CHART_TYPES, render_chart)
R.register_many(XY_CHART_TYPES, render_xy_chart)
