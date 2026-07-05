"""
inspect_pptx.py — 既存の .pptx から「型スペック」を抽出する（②の入力=PowerPoint）。

目的：手元にある良いスライドの構造を、機械的に読み取って JSON 化する。
抽出するのは「どこに・何が・どれだけ」あるか：
  - 各シェイプの種類 / 位置(相対%) / サイズ / テキスト有無 / 塗り色
  - TABLE のセル値（行×列）、CHART の種別・カテゴリ・系列名+数値（実データ）
  - GROUP 図形は再帰的にフラット化し、子図形の座標をスライド絶対%へ変換して展開する
    （group 自体は出力に含めない。中身が見えないと DSL 再構成が骨抜きになるため）。
  - テキストのフォントサイズ階層（見出し/本文/注釈の推定）
  - 配色（使われている色とおおよその面積比 → 70:25:5 になっているか点検）
この JSON を Claude（社内Claude Code）に渡せば、既存型への当てはめ or 新型の設計ができる。

※これは「決定的に読める部分」を担当する。意味の解釈（compare型だ・どのchart型か等）はLLMが行う分業。
　TABLE/CHART は実データをそのまま渡す（LLM が数値を捏造しないよう、根拠を機械抽出側で保証する）。

使い方:
    python -m slidegen.inspect_pptx good_deck.pptx          # 全スライドのスペックをJSON出力
    python -m slidegen.inspect_pptx good_deck.pptx -n 3     # 3枚目だけ
    python -m slidegen.inspect_pptx good_deck.pptx --compact # LLM向けテキストスペック

inspect_compact() はブラウザ（Pyodide worker）のデザイン取り込みが使う LLM 向け出力。
サイズ上限（スライド数・スライドあたり文字数）は関数側で保証する。
"""
from __future__ import annotations
import json, argparse
from collections import Counter
from pptx import Presentation
from pptx.oxml.ns import qn


def _emu_to_pct(v, total):
    try:
        return round(100.0 * float(v) / float(total), 1)
    except Exception:
        return None


def _shape_fill_hex(shape):
    try:
        f = shape.fill
        if f.type is not None and f.fore_color and f.fore_color.type is not None:
            rgb = f.fore_color.rgb
            return str(rgb)
    except Exception:
        pass
    return None


def _font_sizes(shape):
    sizes = []
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size:
                    sizes.append(int(r.font.size.pt))
    return sizes


def _paragraph_text(shape, limit: int) -> str:
    """テキストを取り出す。段落が複数ある（箇条書き等）場合は、レベルを '-' の連続で
    マーカー化して " / " 区切りで連結する（従来は text_frame.text の素通しで階層情報が
    失われ、見出しと箇条書き本文の区別がLLMから付かなかった）。単一段落なら生テキストのみ。
    """
    if not shape.has_text_frame:
        return ""
    paras = list(shape.text_frame.paragraphs)
    texts = [p.text.strip() for p in paras if p.text.strip()]
    if not texts:
        return ""
    if len(texts) == 1:
        return texts[0][:limit]
    parts = []
    for p in paras:
        t = p.text.strip()
        if not t:
            continue
        marker = "-" * ((p.level or 0) + 1) + " "
        parts.append(marker + t)
    return " / ".join(parts)[:limit]


def _table_rows(shape, cell_limit: int = 40, max_rows: int = 12, max_cols: int = 8):
    """表のセル値を行×列で抽出する。

    GraphicFrame は has_text_frame=False のため、従来はテーブルの中身（列見出し・
    セル値）が一切抽出されず「表がある」という事実だけが（テキスト無しの空アイテム
    として）残り、LLM は内容を推測すらできなかった。
    """
    tbl = shape.table
    rows = []
    for r in list(tbl.rows)[:max_rows]:
        rows.append([c.text.strip()[:cell_limit] for c in list(r.cells)[:max_cols]])
    return rows


def _chart_spec(shape, series_limit: int = 6, cat_limit: int = 20):
    """チャートの種別・カテゴリ・系列名+数値（実データ）を抽出する。

    従来はチャートの中身が一切抽出されず、LLM が「それらしい数値」を捏造して
    DSL の bar_chart/line_chart 等を書くリスクがあった。実データを渡すことで
    数値の忠実な引き継ぎを可能にする（型の選択自体はLLMの分業のまま）。
    """
    spec: dict = {"chart_type": str(shape.chart.chart_type).split(" (")[0]}
    try:
        plot = shape.chart.plots[0]
        spec["categories"] = [str(c) for c in plot.categories][:cat_limit]
        series = {}
        for s in list(plot.series)[:series_limit]:
            vals = [round(v, 2) if v is not None else None for v in list(s.values)[:cat_limit]]
            series[s.name or "系列"] = vals
        spec["series"] = series
    except Exception:
        pass  # 種別だけでも「チャートがある」事実は伝わる
    return spec


def _group_child_coord(shape):
    """グループ図形の子座標系(chOff/chExt)を (x, y, cx, cy) で返す。読めなければ None。"""
    try:
        xfrm = shape._element.find(qn("p:grpSpPr")).find(qn("a:xfrm"))
        ch_off = xfrm.find(qn("a:chOff"))
        ch_ext = xfrm.find(qn("a:chExt"))
        return (int(ch_off.get("x")), int(ch_off.get("y")),
                int(ch_ext.get("cx")), int(ch_ext.get("cy")))
    except Exception:
        return None


def _transform(shape, box):
    """box=(abs_left, abs_top, scale_x, scale_y, ch_off_x, ch_off_y) を使い、shape の
    グループ内座標(shape.left/top/width/height)をスライド絶対EMUへ変換する。"""
    abs_left, abs_top, scale_x, scale_y, ch_off_x, ch_off_y = box
    left = shape.left if shape.left is not None else ch_off_x
    top = shape.top if shape.top is not None else ch_off_y
    width = shape.width or 0
    height = shape.height or 0
    return (
        abs_left + (left - ch_off_x) * scale_x,
        abs_top + (top - ch_off_y) * scale_y,
        width * scale_x,
        height * scale_y,
    )


def _iter_leaf_shapes(shapes, box=None):
    """GROUP を再帰的にフラット化し、(shape, abs_left, abs_top, abs_width, abs_height) を
    yield する（group 自身は yield しない）。box は現在の座標系→スライド絶対EMUへの
    変換パラメータ（トップレベルは None＝既に絶対値）。

    従来 group 図形は非再帰で、中身のシェイプ（テキスト・図形）が一切 inspect から
    見えず、group はテキスト無しの空アイテムとして出力されるだけだった。
    """
    for sh in shapes:
        if box is not None:
            abs_left, abs_top, abs_width, abs_height = _transform(sh, box)
        else:
            abs_left, abs_top, abs_width, abs_height = sh.left, sh.top, sh.width, sh.height

        if sh.shape_type is not None and str(sh.shape_type).startswith("GROUP"):
            coord = _group_child_coord(sh)
            if coord is None or not abs_width or not abs_height:
                continue  # 座標が読めないグループは安全に展開できないためスキップ
            ch_off_x, ch_off_y, ch_ext_cx, ch_ext_cy = coord
            scale_x = abs_width / ch_ext_cx if ch_ext_cx else 1.0
            scale_y = abs_height / ch_ext_cy if ch_ext_cy else 1.0
            inner_box = (abs_left, abs_top, scale_x, scale_y, ch_off_x, ch_off_y)
            yield from _iter_leaf_shapes(sh.shapes, box=inner_box)
        else:
            yield sh, abs_left, abs_top, abs_width, abs_height


def inspect_slide(slide, sw, sh, text_limit: int = 40) -> dict:
    shapes = []
    color_area = Counter()
    all_sizes = []
    for sh_, abs_l, abs_t, abs_w, abs_h in _iter_leaf_shapes(slide.shapes):
        item = {
            "kind": str(sh_.shape_type),
            "x%": _emu_to_pct(abs_l, sw) if abs_l is not None else None,
            "y%": _emu_to_pct(abs_t, sh) if abs_t is not None else None,
            "w%": _emu_to_pct(abs_w, sw) if abs_w is not None else None,
            "h%": _emu_to_pct(abs_h, sh) if abs_h is not None else None,
            "text": _paragraph_text(sh_, text_limit),
        }
        if getattr(sh_, "has_table", False):
            item["table"] = _table_rows(sh_)
        if getattr(sh_, "has_chart", False):
            item["chart"] = _chart_spec(sh_)
        fill = _shape_fill_hex(sh_)
        item["fill"] = fill
        sizes = _font_sizes(sh_)
        if sizes:
            item["font_pt"] = sizes
            all_sizes += sizes
        if fill and abs_w and abs_h:
            color_area[fill] += int(abs_w) * int(abs_h)
        shapes.append(item)

    # フォント階層の推定（大→小）
    uniq = sorted(set(all_sizes), reverse=True)
    hierarchy = {}
    if uniq:
        hierarchy["headline_pt"] = uniq[0]
        if len(uniq) > 1:
            hierarchy["body_pt"] = uniq[len(uniq) // 2]
        if len(uniq) > 2:
            hierarchy["caption_pt"] = uniq[-1]

    # 配色の面積比（70:25:5 になっているかの点検材料）
    total_area = sum(color_area.values()) or 1
    palette = [{"color": c, "area%": round(100 * a / total_area, 1)}
               for c, a in color_area.most_common(6)]

    return {
        "n_shapes": len(shapes),
        "font_hierarchy": hierarchy,
        "palette_by_area": palette,
        "shapes": shapes,
    }


def inspect(path: str, only: int | None = None, text_limit: int = 40) -> dict:
    prs = Presentation(path)
    sw, sh = prs.slide_width, prs.slide_height
    out = {"file": path, "slide_size_emu": [int(sw), int(sh)], "slides": []}
    for i, slide in enumerate(prs.slides, 1):
        if only and i != only:
            continue
        spec = inspect_slide(slide, sw, sh, text_limit=text_limit)
        spec["index"] = i
        out["slides"].append(spec)
    return out


def _compact_shape_line(s: dict) -> str:
    kind = s["kind"].split(" (")[0]  # "TEXT_BOX (17)" → "TEXT_BOX"
    geo = ""
    if s["x%"] is not None and s["w%"] is not None:
        geo = f" @({s['x%']},{s['y%']} {s['w%']}x{s['h%']}%)"
    fill = f" fill={s['fill']}" if s.get("fill") else ""
    font = f" {min(s['font_pt'])}-{max(s['font_pt'])}pt" if s.get("font_pt") else ""
    extra = ""
    if s.get("table"):
        # セル区切りは " / "（カンマ区切りにすると "2,580円" のようなセル内容と
        # 列境界が区別できなくなるため）。行区切りは " ; "。
        rows = s["table"]
        extra = " table=" + " ; ".join("[" + " / ".join(row) + "]" for row in rows)
    elif s.get("chart"):
        c = s["chart"]
        cats = " / ".join(c.get("categories", []))
        series = " ".join(
            f"{name}=[{' / '.join('' if v is None else str(v) for v in vals)}]"
            for name, vals in c.get("series", {}).items()
        )
        extra = f" chart={c.get('chart_type', '?')} cats=[{cats}] {series}".rstrip()
    text = f' "{s["text"]}"' if s.get("text") else ""
    return f"  - {kind}{geo}{fill}{font}{extra}{text}"


def inspect_compact(path: str, max_slides: int = 30,
                    chars_per_slide: int = 1600, text_limit: int = 160) -> str:
    """LLM に渡すためのコンパクトなテキストスペックを返す（デザイン取り込み用）。

    - スライド上限・スライドあたり文字数上限を関数側で保証する
      （呼び出し側=ブラウザは gateway の入力上限 200KB を意識しなくてよい）。
    - 内容テキストは text_limit 字まで（DSL 再構成には本文が必要なので JSON 版より長め）。
    - TABLE/CHART は実データ（セル値・カテゴリ・系列数値）を含める（文字数予算は消費するが、
      LLM が数値を捏造しないためには実データの提示が必須）。
    """
    data = inspect(path, text_limit=text_limit)
    lines = [f"deck: {len(data['slides'])} slides"]
    for spec in data["slides"][:max_slides]:
        head = [f"[S{spec['index']}] shapes={spec['n_shapes']}"]
        h = spec["font_hierarchy"]
        if h:
            head.append("fonts(pt): " + " ".join(f"{k.removesuffix('_pt')}={v}" for k, v in h.items()))
        if spec["palette_by_area"]:
            head.append("palette: " + " ".join(f"{p['color']}:{p['area%']}%" for p in spec["palette_by_area"]))
        block = [" ".join(head)]
        used = len(block[0])
        omitted = 0
        for s in spec["shapes"]:
            line = _compact_shape_line(s)
            if used + len(line) > chars_per_slide:
                omitted += 1
                continue
            block.append(line)
            used += len(line) + 1
        if omitted:
            block.append(f"  …({omitted} shapes omitted)")
        lines.append("\n".join(block))
    rest = len(data["slides"]) - max_slides
    if rest > 0:
        lines.append(f"…({rest} more slides omitted)")
    return "\n\n".join(lines)


def main():
    ap = argparse.ArgumentParser(description="既存pptxから型スペックを抽出")
    ap.add_argument("pptx")
    ap.add_argument("-n", type=int, default=None, help="特定スライド番号のみ")
    ap.add_argument("-o", "--output", default=None, help="JSON出力先（省略時stdout）")
    ap.add_argument("--compact", action="store_true", help="LLM向けテキストスペックで出力")
    args = ap.parse_args()
    if args.compact:
        js = inspect_compact(args.pptx)
    else:
        data = inspect(args.pptx, args.n)
        js = json.dumps(data, ensure_ascii=False, indent=2)
    if args.output:
        with open(args.output, "w", encoding="utf-8") as f:
            f.write(js)
        print(f"出力: {args.output}")
    else:
        print(js)


if __name__ == "__main__":
    main()
